"""
NotebookLM integration for generating podcasts, infographics, videos, and
slide decks via notebooklm-py.

This module wraps the notebooklm-py library to create temporary notebooks,
add course content as sources, generate audio/infographic/video/slide-deck
artifacts, and download the results.

Requires NOTEBOOKLM_AUTH_JSON env var with Google session state.
"""
import asyncio
import os
import subprocess
import tempfile
import urllib.request


def _get_event_loop():
    """Get or create an event loop for running async code from sync context."""
    try:
        loop = asyncio.get_event_loop()
        if loop.is_closed():
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
        return loop
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        return loop


NBLM_CHUNK_SIZE = 48000  # safely under the 50,000 char NotebookLM source limit


async def _add_documents_to_notebook(client, notebook_id, documents):
    """Add each document as one or more text sources in the notebook.

    documents: list of dicts with 'filename' and 'extracted_text' keys.
    Documents longer than NBLM_CHUNK_SIZE are split into sequential chunks,
    each added as a separate source so no content is lost.
    """
    for doc in documents:
        text = doc.get('extracted_text', '')
        name = doc.get('filename', 'Document')
        if not text:
            continue
        if len(text) <= NBLM_CHUNK_SIZE:
            await client.sources.add_text(notebook_id, name, text, wait=True, wait_timeout=180.0)
            print(f"[NotebookLM] Added source '{name}' ({len(text)} chars)")
        else:
            chunks = [text[i:i + NBLM_CHUNK_SIZE] for i in range(0, len(text), NBLM_CHUNK_SIZE)]
            for idx, chunk in enumerate(chunks):
                chunk_title = f"{name} (part {idx + 1}/{len(chunks)})"
                await client.sources.add_text(notebook_id, chunk_title, chunk, wait=True, wait_timeout=180.0)
            print(f"[NotebookLM] Added '{name}' as {len(chunks)} chunks ({len(text)} chars total)")


async def _get_or_create_notebook_async(client, title, documents, existing_notebook_id=None):
    """Return the notebook ID to use for generation.

    documents: list of dicts with 'filename' and 'extracted_text'.
    If existing_notebook_id is provided and the notebook is still alive,
    reuse it (adding sources only if none exist yet). Otherwise create a
    fresh notebook and add all documents as sources.
    """
    if existing_notebook_id:
        try:
            await client.notebooks.get(existing_notebook_id)   # raises if deleted
            sources = await client.sources.list(existing_notebook_id)
            if sources:
                print(f"[NotebookLM] Reusing notebook {existing_notebook_id} ({len(sources)} source(s) already present)")
                return existing_notebook_id
            # Notebook exists but sources were never added — add them now
            print(f"[NotebookLM] Reusing notebook {existing_notebook_id}, adding sources")
            await _add_documents_to_notebook(client, existing_notebook_id, documents)
            return existing_notebook_id
        except Exception as e:
            print(f"[NotebookLM] Existing notebook {existing_notebook_id} not usable ({e}), creating new one")

    nb = await client.notebooks.create(title)
    print(f"[NotebookLM] Created new notebook {nb.id}")
    await _add_documents_to_notebook(client, nb.id, documents)
    return nb.id


async def _generate_podcast_async(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Create (or reuse) a NotebookLM notebook, generate podcast, download MP3."""
    from notebooklm import NotebookLMClient
    from notebooklm.rpc.types import AudioFormat, AudioLength

    opts = options or {}
    title = storyboard_json.get('title', 'Course Content')

    # Map string values from frontend to enums
    fmt_map = {
        'DEEP_DIVE': AudioFormat.DEEP_DIVE,
        'BRIEF':     AudioFormat.BRIEF,
        'CRITIQUE':  AudioFormat.CRITIQUE,
        'DEBATE':    AudioFormat.DEBATE,
    }
    len_map = {
        'SHORT':   AudioLength.SHORT,
        'DEFAULT': AudioLength.DEFAULT,
        'LONG':    AudioLength.LONG,
    }
    audio_format = fmt_map.get(opts.get('format', ''), AudioFormat.DEEP_DIVE)
    audio_length = len_map.get(opts.get('length', ''), AudioLength.DEFAULT)

    default_instructions = (
        f"Create an engaging educational podcast about '{title}'. "
        "Make it conversational and suitable for learning. "
        "Cover the key concepts thoroughly."
    )
    instructions = opts.get('instructions') or default_instructions

    async with await NotebookLMClient.from_storage() as client:
        _docs = [{'filename': title, 'extracted_text': source_text}]
        notebook_id = await _get_or_create_notebook_async(client, title, _docs, existing_notebook_id)

        try:
            # Generate audio podcast
            status = await client.artifacts.generate_audio(
                notebook_id,
                instructions=instructions,
                audio_format=audio_format,
                audio_length=audio_length,
            )
            # Allow up to 15 minutes — audio generation is slow
            final = await client.artifacts.wait_for_completion(notebook_id, status.task_id, timeout=900.0)

            # Check if NotebookLM itself reported a failure
            if final.is_failed:
                if final.is_rate_limited:
                    raise Exception('NotebookLM rate limit exceeded — please try again later')
                raise Exception(f'NotebookLM audio generation failed: {final.error or "unknown error"}')

            # Download the generated audio
            await client.artifacts.download_audio(notebook_id, output_path)

            return notebook_id

        except Exception:
            # NOTE: not deleting notebook on failure — keep it for inspection
            raise


async def _generate_infographic_async(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Create (or reuse) a NotebookLM notebook, generate infographic, download PNG."""
    from notebooklm import NotebookLMClient
    from notebooklm.rpc.types import InfographicOrientation, InfographicDetail

    opts = options or {}
    title = storyboard_json.get('title', 'Course Content')

    ori_map = {
        'LANDSCAPE': InfographicOrientation.LANDSCAPE,
        'PORTRAIT':  InfographicOrientation.PORTRAIT,
        'SQUARE':    InfographicOrientation.SQUARE,
    }
    det_map = {
        'CONCISE':  InfographicDetail.CONCISE,
        'STANDARD': InfographicDetail.STANDARD,
        'DETAILED': InfographicDetail.DETAILED,
    }
    # Only pass explicit values when user chose non-default; otherwise pass None
    # so the API uses its own defaults (passing PORTRAIT/STANDARD enums can
    # trigger USER_DISPLAYABLE_ERROR on some account configurations).
    orientation  = ori_map.get(opts.get('orientation', ''))   # None if not specified
    detail_level = det_map.get(opts.get('detail', ''))        # None if not specified
    instructions = opts.get('instructions') or None

    async with await NotebookLMClient.from_storage() as client:
        _docs = [{'filename': title, 'extracted_text': source_text}]
        notebook_id = await _get_or_create_notebook_async(client, title, _docs, existing_notebook_id)

        try:
            # Generate infographic
            print(f"[NotebookLM] Calling generate_infographic: notebook={notebook_id}, orientation={orientation}, detail_level={detail_level}, instructions={bool(instructions)}")
            status = await client.artifacts.generate_infographic(
                notebook_id,
                instructions=instructions,
                orientation=orientation,
                detail_level=detail_level,
            )
            print(f"[NotebookLM] generate_infographic returned status: task_id={getattr(status, 'task_id', None)}, status={getattr(status, 'status', None)}, error={getattr(status, 'error', None)}")
            # Fast-fail: if the API rejected the request immediately, don't poll
            if getattr(status, 'is_failed', False) or not getattr(status, 'task_id', None):
                raise Exception(f'NotebookLM infographic generation rejected: {getattr(status, "error", None) or "no task_id returned"}')
            # Allow up to 15 minutes — infographic generation can be slow
            final = await client.artifacts.wait_for_completion(notebook_id, status.task_id, timeout=900.0)

            # Check if NotebookLM itself reported a failure
            if final.is_failed:
                if final.is_rate_limited:
                    raise Exception('NotebookLM rate limit exceeded — please try again later')
                raise Exception(f'NotebookLM infographic generation failed: {final.error or "unknown error"}')

            # Download the generated infographic
            await client.artifacts.download_infographic(notebook_id, output_path)

            return notebook_id

        except Exception:
            # NOTE: not deleting notebook on failure — keep it for inspection
            raise


async def _generate_video_async(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Create (or reuse) a NotebookLM notebook, generate video, download MP4."""
    from notebooklm import NotebookLMClient
    from notebooklm.rpc.types import VideoFormat, VideoStyle

    opts = options or {}
    title = storyboard_json.get('title', 'Course Content')

    fmt_map = {
        'EXPLAINER': VideoFormat.EXPLAINER,
        'BRIEF':     VideoFormat.BRIEF,
    }
    sty_map = {
        'AUTO_SELECT': VideoStyle.AUTO_SELECT,
        'WHITEBOARD':  VideoStyle.WHITEBOARD,
        'CLASSIC':     VideoStyle.CLASSIC,
        'KAWAII':      VideoStyle.KAWAII,
        'ANIME':       VideoStyle.ANIME,
        'WATERCOLOR':  VideoStyle.WATERCOLOR,
        'RETRO_PRINT': VideoStyle.RETRO_PRINT,
        'HERITAGE':    VideoStyle.HERITAGE,
        'PAPER_CRAFT': VideoStyle.PAPER_CRAFT,
    }
    video_format = fmt_map.get(opts.get('format', ''), VideoFormat.EXPLAINER)
    video_style  = sty_map.get(opts.get('style', ''),  VideoStyle.AUTO_SELECT)

    default_instructions = (
        f"Create an engaging educational video overview of '{title}'. "
        "Make it clear, informative, and suitable for learning. "
        "Cover the key concepts thoroughly."
    )
    instructions = opts.get('instructions') or default_instructions

    async with await NotebookLMClient.from_storage() as client:
        _docs = [{'filename': title, 'extracted_text': source_text}]
        notebook_id = await _get_or_create_notebook_async(client, title, _docs, existing_notebook_id)

        try:
            status = await client.artifacts.generate_video(
                notebook_id,
                instructions=instructions,
                video_format=video_format,
                video_style=video_style,
            )
            # Allow up to 30 minutes — video generation is slower than podcast/infographic
            final = await client.artifacts.wait_for_completion(notebook_id, status.task_id, timeout=1800.0)

            # Check if NotebookLM itself reported a failure
            if final.is_failed:
                if final.is_rate_limited:
                    raise Exception('NotebookLM rate limit exceeded — please try again later')
                raise Exception(f'NotebookLM video generation failed: {final.error or "unknown error"}')

            # Download the generated video
            await client.artifacts.download_video(notebook_id, output_path)

            return notebook_id

        except Exception:
            # NOTE: not deleting notebook on failure — keep it for inspection
            raise


def _pdf_to_images(pdf_path, output_dir):
    """Convert each page of a PDF to a PNG image. Returns list of image file paths.

    Paints a white strip across the bottom ~4 % of each slide to cover
    the NotebookLM branding that is baked into the PDF.
    """
    from pdf2image import convert_from_path
    from PIL import ImageDraw

    images = convert_from_path(pdf_path, dpi=200, fmt='png')
    paths = []
    for i, img in enumerate(images):
        # Cover NotebookLM logo / branding bar at bottom of slide
        draw = ImageDraw.Draw(img)
        strip_h = max(int(img.height * 0.04), 20)
        draw.rectangle([0, img.height - strip_h, img.width, img.height], fill='white')

        img_path = os.path.join(output_dir, f'slide_{i+1}.png')
        img.save(img_path, 'PNG')
        paths.append(img_path)
    return paths


def _extract_shape_texts(shape):
    """Recursively extract all text from a shape, including tables and groups."""
    texts = []
    # Text frames (titles, body text, callouts, etc.)
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                texts.append(text)
    # Tables
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                ct = cell.text.strip()
                if ct:
                    row_texts.append(ct)
            if row_texts:
                texts.append(' | '.join(row_texts))
    # Group shapes — recurse into children
    if shape.shape_type is not None:
        try:
            for child in shape.shapes:
                texts.extend(_extract_shape_texts(child))
        except (AttributeError, TypeError):
            pass
    return texts


def _extract_slide_texts(pptx_path):
    """Extract text content from each slide in a PPTX file.

    Handles text frames, tables, and group shapes so voiceover prompts
    get the full content of every slide.
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            texts.extend(_extract_shape_texts(shape))
        slide_texts.append('\n'.join(texts))
    return slide_texts


def _clean_voiceover_script(text):
    """Remove reference numbers, citations, and other TTS-unfriendly artifacts."""
    import re
    # Bracketed references: [1], [2, 3], [ref], [i], [ii]
    text = re.sub(r'\[\d+(?:\s*,\s*\d+)*\]', '', text)
    text = re.sub(r'\[[ivxlc]+\]', '', text, flags=re.IGNORECASE)
    # Parenthetical citations: (p. 42), (pp. 10-15), (Smith, 2024), (Section 3.2)
    text = re.sub(r'\(pp?\.\s*\d[\d\-–, ]*\)', '', text)
    text = re.sub(r'\([A-Z][a-z]+(?:\s+(?:et\s+al\.?|&|and)\s+[A-Z][a-z]+)?,?\s*\d{4}\)', '', text)
    text = re.sub(r'\((?:Section|Clause|Article|Chapter|Part|Appendix|Table|Figure)\s+[\d.]+[a-z]?\)', '', text, flags=re.IGNORECASE)
    # Standard codes: AS/NZS 4801:2001, ISO 9001:2015, BS EN 1234
    text = re.sub(r'\b[A-Z]{2,}(?:/[A-Z]{2,})*\s+\d{3,}(?::\d{4})?\b', '', text)
    # Section/clause references inline: Section 3.2.1, Clause 4.5, Article 12(b)
    text = re.sub(r'\b(?:Section|Clause|Article|Appendix)\s+\d+(?:\.\d+)*(?:\([a-z]\))?', '', text, flags=re.IGNORECASE)
    # Document reference codes: DOC-2024-0451, REF-123, ID: ABC-123
    text = re.sub(r'\b[A-Z]{2,}-\d{2,}(?:-\d+)*\b', '', text)
    text = re.sub(r'\b(?:Ref|ID|Doc)[\s.:]+[A-Z0-9][\w\-]*', '', text, flags=re.IGNORECASE)
    # Footnote markers
    text = re.sub(r'[*†‡§¶]+(?=\s|$)', '', text)
    # Markdown bold/italic
    text = re.sub(r'\*{1,2}(.+?)\*{1,2}', r'\1', text)
    # Clean up leftover punctuation artifacts and multiple spaces
    text = re.sub(r'\(\s*\)', '', text)        # empty parens
    text = re.sub(r'\[\s*\]', '', text)        # empty brackets
    text = re.sub(r'\s{2,}', ' ', text)        # collapse spaces
    text = re.sub(r'\s+([.,;:!?])', r'\1', text)  # space before punctuation
    return text.strip()


def _parse_voiceover_scripts(response_text, expected_count):
    """Parse [SLIDE N] markers from NBLM response into a list of per-slide scripts.

    Matches by slide number (not occurrence order) so out-of-order responses
    from NBLM are still correctly mapped to the right slide.
    """
    import re
    pattern = re.compile(r'\[SLIDE\s+(\d+)\](.*?)(?=\[SLIDE\s+\d+\]|$)', re.DOTALL)
    scripts_by_num = {}
    for num_str, content in pattern.findall(response_text):
        num = int(num_str)
        cleaned = _clean_voiceover_script(content.strip())
        if cleaned:
            scripts_by_num[num] = cleaned
    # Build ordered list indexed 0 … expected_count-1
    return [scripts_by_num.get(i + 1, '') for i in range(expected_count)]


THEME_INSTRUCTIONS = {
    'corporate': 'Style the slides with a white background, deep navy headings (#1e3a5f), dark charcoal body text, and blue accents (#2c5aa0). Use clean professional fonts such as Calibri or Arial.',
    'dark':      'Style the slides with a very dark background (#12121f), white headings, light blue-grey body text (#c8d2e6), and bright blue accents (#4facfe). Use modern sans-serif fonts.',
    'slate':     'Style the slides with a light grey background (#f5f6f8), dark charcoal headings and body text (#2d3436), and blue accents (#639fff). Use clean minimal fonts.',
    'clean':     'Style the slides with a white background, near-black headings (#1a1a2a), dark body text, and green accents (#30d158). Use crisp modern fonts.',
    'vibrant':   'Style the slides with a white background, ocean blue headings (#0077b6), dark body text, and orange accents (#ff6b35). Use bold energetic fonts.',
}


async def _apply_theme_to_notebook(opts, client, notebook_id, tmp_dir):
    """Apply presentation theme styling to a NotebookLM notebook.

    For predefined themes: returns a style instruction string (no source added).
    For custom .pptx templates: converts to PDF via LibreOffice, uploads as NBLM
    source, and returns a matching instruction string.
    Falls back to text extraction if LibreOffice conversion fails.
    """
    ppt_theme = opts.get('ppt_theme')
    ppt_template_url = opts.get('ppt_template_url')

    if ppt_template_url:
        try:
            tmp_pptx = os.path.join(tmp_dir, 'template.pptx')
            print(f'[theme] Downloading template from {ppt_template_url}')
            await asyncio.to_thread(urllib.request.urlretrieve, ppt_template_url, tmp_pptx)

            # Convert to PDF via LibreOffice headless
            print('[theme] Converting template to PDF via LibreOffice')
            await asyncio.to_thread(
                lambda: subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', tmp_dir, tmp_pptx],
                    timeout=60,
                    check=True,
                    capture_output=True,
                )
            )

            pdf_path = os.path.join(tmp_dir, 'template.pdf')
            if not os.path.exists(pdf_path):
                raise FileNotFoundError('LibreOffice did not produce a PDF')

            print('[theme] Uploading template PDF as NBLM source')
            with open(pdf_path, 'rb') as f:
                await client.sources.add_file(
                    notebook_id, f,
                    mime_type='application/pdf',
                    title='Template Style Reference',
                    wait=True,
                )
            return 'Match the visual style of the uploaded template (added as a source). Replicate its color scheme, fonts, and layout structure.'

        except Exception as e:
            print(f'[theme] LibreOffice conversion failed ({e}), falling back to text extraction')
            try:
                from pptx import Presentation
                prs = Presentation(tmp_pptx)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            texts.append(shape.text.strip())
                template_text = '\n'.join(texts[:200])  # cap at 200 snippets
                if template_text:
                    await client.sources.add_text(
                        notebook_id,
                        'Template Style Reference',
                        template_text[:50000],
                        wait=True,
                    )
                return 'Use the visual style suggested by the uploaded template reference text. Match its tone, structure, and implied color scheme.'
            except Exception as e2:
                print(f'[theme] Text fallback also failed ({e2}), skipping template styling')
                return ''

    if ppt_theme and ppt_theme in THEME_INSTRUCTIONS:
        print(f'[theme] Applying predefined theme: {ppt_theme}')
        return THEME_INSTRUCTIONS[ppt_theme]

    return ''


async def _generate_slide_deck_async(documents, title, output_dir, options=None, existing_notebook_id=None, notebook_id_callback=None):
    """Create (or reuse) a NotebookLM notebook, generate slide deck, download PDF + PPTX + per-slide PNGs."""
    from notebooklm import NotebookLMClient
    from notebooklm.rpc.types import SlideDeckFormat, SlideDeckLength

    opts = options or {}

    fmt_map = {
        'DETAILED_DECK':    SlideDeckFormat.DETAILED_DECK,
        'PRESENTER_SLIDES': SlideDeckFormat.PRESENTER_SLIDES,
    }
    len_map = {
        'DEFAULT': SlideDeckLength.DEFAULT,
        'SHORT':   SlideDeckLength.SHORT,
    }
    slide_format = fmt_map.get(opts.get('slide_format', ''))
    slide_length = len_map.get(opts.get('slide_length', ''))
    instructions = opts.get('instructions') or None

    # Inject slide count into instructions — NBLM has no direct count parameter
    slide_count = opts.get('slide_count')
    if slide_count:
        count_instr = f"Create exactly {slide_count} slides."
        instructions = (instructions + '\n\n' + count_instr) if instructions else count_instr

    async with await NotebookLMClient.from_storage() as client:
        notebook_id = await _get_or_create_notebook_async(client, title, documents, existing_notebook_id)

        # Persist the notebook_id immediately so it can be used to retry if we time out later
        if notebook_id_callback:
            try:
                notebook_id_callback(notebook_id)
            except Exception as cb_err:
                print(f"[NotebookLM] notebook_id_callback failed (non-fatal): {cb_err}")

        # Apply theme styling
        if opts.get('ppt_theme') or opts.get('ppt_template_url'):
            try:
                theme_instr = await _apply_theme_to_notebook(opts, client, notebook_id, output_dir)
                if theme_instr:
                    instructions = (instructions + '\n\n' + theme_instr) if instructions else theme_instr
            except Exception as theme_err:
                print(f"[theme] Failed to apply theme (non-fatal): {theme_err}")

        try:
            print(f"[NotebookLM] Generating slide deck: notebook={notebook_id}, format={slide_format}, length={slide_length}")
            status = await client.artifacts.generate_slide_deck(
                notebook_id,
                instructions=instructions,
                slide_format=slide_format,
                slide_length=slide_length,
            )
            if getattr(status, 'is_failed', False) or not getattr(status, 'task_id', None):
                raise Exception(f'NotebookLM slide deck generation rejected: {getattr(status, "error", None) or "no task_id returned"}')

            # Allow up to 60 minutes
            final = await client.artifacts.wait_for_completion(notebook_id, status.task_id, timeout=3600.0)

            if final.is_failed:
                if final.is_rate_limited:
                    raise Exception('NotebookLM rate limit exceeded — please try again later')
                raise Exception(f'NotebookLM slide deck generation failed: {final.error or "unknown error"}')

            # Download as PDF (for preview images)
            pdf_path = os.path.join(output_dir, 'slides.pdf')
            await client.artifacts.download_slide_deck(notebook_id, pdf_path)

            # Download as PPTX (for final download)
            pptx_path = os.path.join(output_dir, 'slides.pptx')
            try:
                await client.artifacts.download_slide_deck(notebook_id, pptx_path, output_format='pptx')
            except Exception as e:
                print(f"[NotebookLM] PPTX download failed ({e}), will use PDF only")
                pptx_path = None

            # Convert PDF pages to individual PNGs for preview
            slide_image_paths = _pdf_to_images(pdf_path, output_dir)
            print(f"[NotebookLM] Slide deck generated: {len(slide_image_paths)} slides")

            # Generate voiceover scripts via NBLM chat
            voiceover_scripts = []
            try:
                # Extract text from PPTX if available; fall back to generic labels
                slide_texts = None
                if pptx_path and os.path.exists(pptx_path):
                    slide_texts = _extract_slide_texts(pptx_path)
                    # Verify extraction returned meaningful content
                    non_empty = sum(1 for t in slide_texts if t.strip())
                    if non_empty == 0:
                        print(f"[NotebookLM] PPTX text extraction returned all empty slides, falling back")
                        slide_texts = None
                    else:
                        print(f"[NotebookLM] Extracted text from {non_empty}/{len(slide_texts)} slides")

                if slide_texts is None:
                    slide_texts = [f"(visual content)" for _ in range(len(slide_image_paths))]

                prompt_parts = []
                for i, text in enumerate(slide_texts):
                    label = text.strip() if text.strip() else "(visual/graphical content — no extractable text)"
                    prompt_parts.append(f"[SLIDE {i+1}]\n{label}")

                # Build timing constraint if provided
                target_time = opts.get('target_time') if opts else None
                max_time = opts.get('max_time') if opts else None
                timing_instruction = ""
                if target_time and max_time:
                    avg_per_slide = round(target_time * 60 / len(slide_texts))
                    max_per_slide = round(max_time * 60 / len(slide_texts))
                    timing_instruction = (
                        f"\n\nIMPORTANT TIMING CONSTRAINT: The total voiceover narration for all slides combined "
                        f"should target approximately {target_time} minutes and must not exceed {max_time} minutes. "
                        f"With {len(slide_texts)} slides, aim for roughly {avg_per_slide} seconds per slide "
                        f"(maximum {max_per_slide} seconds per slide). "
                        f"A typical speaking pace is about 150 words per minute. "
                        f"Adjust script length per slide accordingly — some slides may need shorter scripts, "
                        f"others longer, but the total should stay within the time budget."
                    )
                elif target_time:
                    avg_per_slide = round(target_time * 60 / len(slide_texts))
                    timing_instruction = (
                        f"\n\nTIMING GUIDELINE: The total voiceover narration should target approximately "
                        f"{target_time} minutes ({avg_per_slide} seconds average per slide at ~150 words/minute)."
                    )

                prompt = (
                    f"You have created a {len(slide_texts)}-slide presentation based on the uploaded source documents.\n"
                    "Write a voiceover narration script that a presenter will read aloud while showing each slide.\n\n"
                    "CRITICAL RULE: Each slide's narration must ONLY describe the topics and content that appear on THAT SPECIFIC slide. "
                    "Do NOT reorganise, merge, or move content between slides. "
                    "If a slide's text mentions Topic X, the narration for that slide must cover Topic X — not a different topic from another slide.\n\n"
                    "Here is the exact text content extracted from each slide:\n\n"
                    + '\n\n'.join(prompt_parts) +
                    "\n\nRespond with each slide's script preceded by [SLIDE N] on its own line:\n\n"
                    "[SLIDE 1]\n(narration for slide 1)\n\n"
                    "[SLIDE 2]\n(narration for slide 2)\n\n"
                    f"... and so on for all {len(slide_texts)} slides.\n\n"
                    "Guidelines:\n"
                    "- Stay faithful to each slide's content — narrate what is ON the slide, enriched with background from the source documents\n"
                    "- Keep the same order as the slides — do not reorder topics\n"
                    "- For slides with only visual/graphical content, describe the visual and its meaning based on the source material\n"
                    "- Make the narration natural and conversational, suitable for a professional training presentation\n\n"
                    "IMPORTANT: This script will be read aloud by a text-to-speech system. "
                    "Do NOT include document reference numbers, citation markers, section numbers, "
                    "standard codes (e.g. 'AS/NZS 4801'), page references, footnote markers, "
                    "or any alphanumeric identifiers. Paraphrase content naturally instead of citing sources."
                    + timing_instruction
                )

                print(f"[NotebookLM] Generating voiceover scripts for {len(slide_texts)} slides...")
                result = await client.chat.ask(notebook_id, prompt)
                answer_text = getattr(result, 'answer', '') or str(result)
                voiceover_scripts = _parse_voiceover_scripts(answer_text, len(slide_image_paths))
                print(f"[NotebookLM] Voiceover scripts generated: {len(voiceover_scripts)} scripts")
            except Exception as e:
                print(f"[NotebookLM] Voiceover script generation failed ({e}), returning empty scripts")
                voiceover_scripts = [''] * len(slide_image_paths)

            # Add voiceover scripts back to the notebook as a source
            if voiceover_scripts:
                try:
                    scripts_text = '\n\n'.join(
                        f"[SLIDE {i + 1}]\n{script}"
                        for i, script in enumerate(voiceover_scripts)
                        if script and script.strip()
                    )
                    if scripts_text:
                        await client.sources.add_text(
                            notebook_id,
                            'Voiceover Scripts',
                            scripts_text,
                            wait=True,
                            wait_timeout=180.0,
                        )
                        print(f"[NotebookLM] Added voiceover scripts as source ({len(scripts_text)} chars)")
                except Exception as e:
                    print(f"[NotebookLM] Failed to add voiceover scripts as source (non-fatal): {e}")

            return notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts

        except Exception:
            raise


async def _cleanup_notebook_async(notebook_id):
    """Delete a temporary NotebookLM notebook."""
    from notebooklm import NotebookLMClient

    async with await NotebookLMClient.from_storage() as client:
        await client.notebooks.delete(notebook_id)


async def _check_auth_async():
    """Test NotebookLM auth by listing notebooks."""
    from notebooklm import NotebookLMClient

    async with await NotebookLMClient.from_storage() as client:
        await client.notebooks.list()
    return True


# ── Public sync wrappers ──────────────────────────────────────────────────────

def generate_podcast(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Sync wrapper: generate a podcast MP3 from course content."""
    loop = _get_event_loop()
    return loop.run_until_complete(
        _generate_podcast_async(source_text, storyboard_json, output_path, options, existing_notebook_id)
    )


def generate_infographic(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Sync wrapper: generate an infographic PNG from course content."""
    loop = _get_event_loop()
    return loop.run_until_complete(
        _generate_infographic_async(source_text, storyboard_json, output_path, options, existing_notebook_id)
    )


def generate_video(source_text, storyboard_json, output_path, options=None, existing_notebook_id=None):
    """Sync wrapper: generate a video MP4 from course content."""
    loop = _get_event_loop()
    return loop.run_until_complete(
        _generate_video_async(source_text, storyboard_json, output_path, options, existing_notebook_id)
    )


def generate_slide_deck(documents, title, output_dir, options=None, existing_notebook_id=None, notebook_id_callback=None):
    """Sync wrapper: generate a slide deck PDF + PPTX + preview images from content.

    documents: list of dicts with 'filename' and 'extracted_text' keys.
    """
    loop = _get_event_loop()
    return loop.run_until_complete(
        _generate_slide_deck_async(documents, title, output_dir, options, existing_notebook_id, notebook_id_callback)
    )


def cleanup_notebook(notebook_id):
    """Sync wrapper: delete a temporary NotebookLM notebook."""
    loop = _get_event_loop()
    loop.run_until_complete(_cleanup_notebook_async(notebook_id))


def check_auth():
    """Sync wrapper: test whether NotebookLM auth is valid. Returns True or raises."""
    loop = _get_event_loop()
    return loop.run_until_complete(_check_auth_async())
