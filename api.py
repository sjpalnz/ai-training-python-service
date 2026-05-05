from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import json
import os
import threading
from datetime import datetime
from generate_powerpoint import generate_powerpoint_file
from generate_scorm import generate_scorm_package


# ── Google Drive helpers ──────────────────────────────────────────────────────

def get_drive_credentials(refresh_token: str):
    """Exchange a refresh token for a valid Google Credentials object."""
    from google.auth.transport.requests import Request
    from google.oauth2.credentials import Credentials
    creds = Credentials(
        token=None,
        refresh_token=refresh_token,
        token_uri='https://oauth2.googleapis.com/token',
        client_id=os.environ.get('GOOGLE_OAUTH_CLIENT_ID'),
        client_secret=os.environ.get('GOOGLE_OAUTH_CLIENT_SECRET')
    )
    creds.refresh(Request())
    return creds


def build_drive_service(refresh_token: str):
    """Return an authenticated Google Drive v3 service."""
    from googleapiclient.discovery import build
    return build('drive', 'v3', credentials=get_drive_credentials(refresh_token))

app = Flask(__name__)
CORS(app)  # Allow requests from Supabase

# Directory for generated files
OUTPUT_DIR = '/tmp/generated_files'
os.makedirs(OUTPUT_DIR, exist_ok=True)


def _embed_in_background(doc_id: str, client_id: str, filename: str, extracted_text: str):
    """Run RAG chunking + embedding in a background thread so the HTTP response
    is not blocked by the (potentially slow) embedding step."""
    try:
        from embeddings import embed_texts, chunk_text
        supabase = get_supabase_client()
        supabase.table('document_chunks').delete().eq('document_id', doc_id).execute()
        chunks = chunk_text(extracted_text)
        embeddings = embed_texts(chunks)
        chunk_rows = [
            {'document_id': doc_id, 'client_id': client_id, 'chunk_index': i,
             'chunk_text': c, 'embedding': e}
            for i, (c, e) in enumerate(zip(chunks, embeddings))
        ]
        for batch_start in range(0, len(chunk_rows), 100):
            supabase.table('document_chunks').insert(chunk_rows[batch_start:batch_start + 100]).execute()
        print(f"[RAG] {len(chunks)} chunks indexed for {filename}")
    except Exception as embed_err:
        print(f"[RAG] Warning: could not index chunks for {filename}: {embed_err}")


def _prewarm_embedding_model():
    """Load the fastembed model at startup so the first import request is fast."""
    try:
        from embeddings import get_model
        get_model()
        print("[startup] Embedding model pre-loaded and ready.")
    except Exception as e:
        print(f"[startup] Warning: could not pre-load embedding model: {e}")

threading.Thread(target=_prewarm_embedding_model, daemon=True).start()

def verify_jwt(req):
    """Verify a Supabase-issued JWT via the Supabase auth API and return the user UUID, or None if invalid."""
    auth_header = req.headers.get('Authorization', '')
    if not auth_header.startswith('Bearer '):
        return None
    token = auth_header[7:]
    try:
        supabase = get_supabase_client()
        response = supabase.auth.get_user(token)
        return response.user.id
    except Exception:
        return None


def get_supabase_client():
    """Create and return a Supabase client using environment variables"""
    from supabase import create_client
    supabase_url = os.environ.get('SUPABASE_URL')
    supabase_key = os.environ.get('SUPABASE_SERVICE_ROLE_KEY')
    if not supabase_url or not supabase_key:
        raise Exception('SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY environment variables required')
    return create_client(supabase_url, supabase_key)

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'service': 'AI Training Platform - Python Service',
        'timestamp': datetime.now().isoformat()
    })

# --- Extract text from uploaded PowerPoint ---

@app.route('/extract-pptx-text', methods=['POST'])
def extract_pptx_text():
    """Extract structured slide text from an uploaded .pptx file."""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'No file provided'}), 400

        file = request.files['file']
        if not file.filename.lower().endswith('.pptx'):
            return jsonify({'success': False, 'error': 'File must be a .pptx file'}), 400

        from pptx import Presentation
        import io

        prs = Presentation(io.BytesIO(file.read()))
        slides = []

        for idx, slide in enumerate(prs.slides, 1):
            title_text = ''
            body_text_parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    # First text shape with content is usually the title
                    if not title_text and (shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False):
                        title_text = text
                    elif not title_text and idx == 1:
                        title_text = text
                    else:
                        body_text_parts.append(text)

            # Fallback: if no title found, use first body text
            if not title_text and body_text_parts:
                title_text = body_text_parts.pop(0)

            slides.append({
                'number': idx,
                'title': title_text,
                'text': '\n'.join(body_text_parts)
            })

        return jsonify({
            'success': True,
            'slides': slides,
            'filename': file.filename,
            'slide_count': len(slides)
        })

    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


# --- Supabase-integrated endpoints (called by Supabase Edge Functions) ---

@app.route('/generate-powerpoint-from-storyboard', methods=['POST'])
def generate_ppt_from_storyboard():
    """
    Fetch storyboard from Supabase, generate real PowerPoint, upload to storage, save to DB.
    Expected JSON: { "storyboard_id": "uuid", "theme_id": "corporate", "template_url": null }
    """
    import urllib.request
    template_tmpfile = None
    infographic_tmpfile = None
    try:
        data = request.json
        storyboard_id = data.get('storyboard_id')
        theme_id      = data.get('theme_id', 'corporate')
        template_url  = data.get('template_url')   # optional user .pptx template

        if not storyboard_id:
            return jsonify({'error': 'storyboard_id is required'}), 400

        supabase = get_supabase_client()

        # Fetch storyboard and related course
        result = supabase.table('storyboards').select('*, courses(*)').eq('id', storyboard_id).single().execute()
        storyboard = result.data

        if not storyboard:
            return jsonify({'error': 'Storyboard not found'}), 404

        course_data = storyboard['content_json']

        # Download user template if provided
        template_path = None
        if template_url:
            template_tmpfile = os.path.join(OUTPUT_DIR, f"tpl_{storyboard_id[:8]}.pptx")
            urllib.request.urlretrieve(template_url, template_tmpfile)
            template_path = template_tmpfile

        # Check for existing infographic to embed
        infographic_tmpfile = None
        try:
            infographic_result = supabase.table('generated_files') \
                .select('file_url') \
                .eq('course_id', storyboard['course_id']) \
                .eq('file_type', 'infographic') \
                .order('created_at', desc=True) \
                .limit(1) \
                .execute()
            if infographic_result.data:
                infographic_tmpfile = os.path.join(OUTPUT_DIR, f"infographic_{storyboard_id[:8]}.png")
                urllib.request.urlretrieve(infographic_result.data[0]['file_url'], infographic_tmpfile)
        except Exception as e:
            print(f"[PPT] Warning: could not fetch infographic for embedding: {e}")
            infographic_tmpfile = None

        # Generate PowerPoint file
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"course_{storyboard_id[:8]}_{timestamp}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        generate_powerpoint_file(course_data, filepath, theme_id=theme_id, template_path=template_path, infographic_path=infographic_tmpfile)

        # Upload to Supabase Storage
        storage_path = f"powerpoints/{filename}"
        with open(filepath, 'rb') as f:
            file_bytes = f.read()

        supabase.storage.from_('course-files').upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
        )

        # Get public URL
        public_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Save to generated_files table
        file_record = supabase.table('generated_files').insert({
            'course_id': storyboard['course_id'],
            'file_type': 'powerpoint',
            'file_url': public_url,
            'file_size': len(file_bytes)
        }).execute()

        # Cleanup local files
        os.remove(filepath)
        if template_tmpfile and os.path.exists(template_tmpfile):
            os.remove(template_tmpfile)
        if infographic_tmpfile and os.path.exists(infographic_tmpfile):
            os.remove(infographic_tmpfile)

        return jsonify({
            'success': True,
            'file_id': file_record.data[0]['id'],
            'download_url': public_url,
            'filename': filename
        })

    except Exception as e:
        if template_tmpfile and os.path.exists(template_tmpfile):
            os.remove(template_tmpfile)
        if infographic_tmpfile and os.path.exists(infographic_tmpfile):
            os.remove(infographic_tmpfile)
        return jsonify({'error': str(e)}), 500


@app.route('/generate-scorm-from-storyboard', methods=['POST'])
def generate_scorm_from_storyboard():
    """
    Fetch storyboard from Supabase, generate SCORM package, upload to storage, save to DB.
    Expected JSON: { "storyboard_id": "uuid" }
    """
    try:
        data = request.json
        storyboard_id = data.get('storyboard_id')

        if not storyboard_id:
            return jsonify({'error': 'storyboard_id is required'}), 400

        supabase = get_supabase_client()

        # Fetch storyboard and related course
        result = supabase.table('storyboards').select('*, courses(*)').eq('id', storyboard_id).single().execute()
        storyboard = result.data

        if not storyboard:
            return jsonify({'error': 'Storyboard not found'}), 404

        course_data = storyboard['content_json']

        # Check for existing podcast and infographic to embed
        podcast_url = None
        infographic_url = None
        try:
            podcast_result = supabase.table('generated_files') \
                .select('file_url') \
                .eq('course_id', storyboard['course_id']) \
                .eq('file_type', 'podcast') \
                .order('created_at', desc=True) \
                .limit(1) \
                .execute()
            if podcast_result.data:
                podcast_url = podcast_result.data[0]['file_url']

            infographic_result = supabase.table('generated_files') \
                .select('file_url') \
                .eq('course_id', storyboard['course_id']) \
                .eq('file_type', 'infographic') \
                .order('created_at', desc=True) \
                .limit(1) \
                .execute()
            if infographic_result.data:
                infographic_url = infographic_result.data[0]['file_url']
        except Exception as e:
            print(f"[SCORM] Warning: could not fetch media for embedding: {e}")

        # Generate SCORM package
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"scorm_{storyboard_id[:8]}_{timestamp}.zip"
        filepath = os.path.join(OUTPUT_DIR, filename)
        generate_scorm_package(course_data, filepath, podcast_url=podcast_url, infographic_url=infographic_url)

        # Upload to Supabase Storage
        storage_path = f"scorm/{filename}"
        with open(filepath, 'rb') as f:
            file_bytes = f.read()

        supabase.storage.from_('course-files').upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": "application/zip"}
        )

        # Get public URL
        public_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Save to generated_files table
        file_record = supabase.table('generated_files').insert({
            'course_id': storyboard['course_id'],
            'file_type': 'scorm',
            'file_url': public_url,
            'file_size': len(file_bytes)
        }).execute()

        # Cleanup local file
        os.remove(filepath)

        return jsonify({
            'success': True,
            'file_id': file_record.data[0]['id'],
            'download_url': public_url,
            'filename': filename
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Google Drive endpoints ---

@app.route('/check-google-connection', methods=['POST'])
def check_google_connection():
    """
    Check whether a user has Google Drive credentials stored.
    Expects JSON: { user_id }
    Returns: { connected: bool }
    """
    try:
        data = request.json or {}
        user_id = data.get('user_id')
        if not user_id:
            return jsonify({'connected': False}), 200
        supabase = get_supabase_client()
        creds_row = supabase.table('user_google_credentials') \
            .select('id') \
            .eq('user_id', user_id) \
            .maybeSingle() \
            .execute()
        return jsonify({'connected': bool(creds_row.data)})
    except Exception as e:
        return jsonify({'connected': False, 'error': str(e)}), 200


@app.route('/list-google-drive-files', methods=['POST'])
def list_google_drive_files():
    """
    List folders + supported files inside a specific Drive folder.
    Expects JSON: { user_id, folder_id? }  (folder_id defaults to 'root')
    Returns: { success, folders: [{id, name}], files: [{id, name, mimeType, size, modifiedTime}] }
    """
    try:
        data = request.json or {}
        user_id = data.get('user_id')
        folder_id = data.get('folder_id', 'root')

        if not user_id:
            return jsonify({'error': 'user_id required'}), 400

        # Fetch refresh_token from Supabase using service_role (bypasses RLS)
        supabase = get_supabase_client()
        creds_row = supabase.table('user_google_credentials') \
            .select('refresh_token') \
            .eq('user_id', user_id) \
            .single() \
            .execute()
        refresh_token = creds_row.data.get('refresh_token') if creds_row.data else None
        if not refresh_token:
            return jsonify({'error': 'Google Drive not connected for this user'}), 400

        drive_service = build_drive_service(refresh_token)

        common_args = dict(
            spaces='drive',
            pageSize=100,
            orderBy='name',
            includeItemsFromAllDrives=True,
            supportsAllDrives=True
        )

        # Fetch subfolders in this location
        folder_results = drive_service.files().list(
            q=f"'{folder_id}' in parents and mimeType = 'application/vnd.google-apps.folder' and trashed = false",
            fields='files(id, name)',
            **common_args
        ).execute()

        # Fetch supported files in this location
        file_query = (
            f"'{folder_id}' in parents and "
            "(mimeType = 'application/pdf' or "
            "mimeType = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or "
            "mimeType = 'application/vnd.google-apps.document' or "
            "mimeType = 'text/plain') and "
            "trashed = false"
        )
        file_results = drive_service.files().list(
            q=file_query,
            fields='files(id, name, mimeType, size, modifiedTime)',
            **common_args
        ).execute()

        return jsonify({
            'success': True,
            'folders': folder_results.get('files', []),
            'files': file_results.get('files', [])
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/check-drive-updates', methods=['POST'])
def check_drive_updates():
    """Compare stored drive_modified_time with current Drive modifiedTime for each imported doc."""
    user_id = verify_jwt(request)
    if not user_id:
        return jsonify({'error': 'Unauthorized'}), 401

    try:
        supabase = get_supabase_client()

        # Get all Drive-imported docs that have a stored modified time
        docs_result = supabase.table('documents') \
            .select('id, drive_file_id, drive_modified_time') \
            .eq('client_id', user_id) \
            .eq('source', 'google_drive') \
            .not_.is_('drive_file_id', 'null') \
            .not_.is_('drive_modified_time', 'null') \
            .execute()

        drive_docs = docs_result.data or []
        if not drive_docs:
            return jsonify({'updated_docs': []})

        # Get refresh token
        creds_result = supabase.table('user_google_credentials') \
            .select('refresh_token') \
            .eq('user_id', user_id) \
            .limit(1) \
            .execute()

        if not creds_result.data or not creds_result.data[0].get('refresh_token'):
            return jsonify({'updated_docs': []})

        refresh_token = creds_result.data[0]['refresh_token']
        drive_service = build_drive_service(refresh_token)

        updated = []
        for doc in drive_docs:
            try:
                meta = drive_service.files().get(
                    fileId=doc['drive_file_id'],
                    fields='modifiedTime'
                ).execute()
                current_modified = meta.get('modifiedTime', '')
                stored_modified = doc.get('drive_modified_time', '')
                if current_modified and stored_modified and current_modified > stored_modified:
                    updated.append({
                        'doc_id': doc['id'],
                        'drive_file_id': doc['drive_file_id'],
                        'current_modified_time': current_modified,
                        'imported_modified_time': stored_modified
                    })
            except Exception:
                pass  # File may have been deleted from Drive

        return jsonify({'updated_docs': updated})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Document upload and text extraction ---

@app.route('/process-documents', methods=['POST'])
def process_documents():
    """
    Accept uploaded files OR a Google Drive file ID, extract text, save to Supabase.
    Manual upload: multipart form data with 'files' field.
    Drive import:  JSON body with 'google_drive_file_id' and 'refresh_token'.
    Supports: PDF, DOCX, TXT (+ Google Docs exported as DOCX)
    """
    try:
        user_id = verify_jwt(request)
        if not user_id:
            return jsonify({'error': 'Unauthorized'}), 401

        client_id = user_id

        # ── Google Drive import path ──────────────────────────────────────────
        if request.is_json:
            data = request.json or {}
            google_drive_file_id = data.get('google_drive_file_id')

            if not google_drive_file_id:
                return jsonify({'error': 'google_drive_file_id required for Drive import'}), 400

            # Fetch refresh_token from Supabase (service_role bypasses RLS)
            supabase = get_supabase_client()
            creds_row = supabase.table('user_google_credentials') \
                .select('refresh_token') \
                .eq('user_id', user_id) \
                .single() \
                .execute()
            refresh_token = creds_row.data.get('refresh_token') if creds_row.data else None
            if not refresh_token:
                return jsonify({'error': 'Google Drive not connected for this user'}), 400

            drive_service = build_drive_service(refresh_token)

            # Fetch file metadata
            file_meta = drive_service.files().get(
                fileId=google_drive_file_id,
                fields='name, mimeType, size, modifiedTime'
            ).execute()

            filename = file_meta.get('name', f'drive_{google_drive_file_id}')
            mime_type = file_meta.get('mimeType', '')
            file_size = int(file_meta.get('size', 0))
            drive_modified_time = file_meta.get('modifiedTime')  # ISO 8601 string

            # Determine file extension + download method
            from io import BytesIO
            if mime_type == 'application/vnd.google-apps.document':
                # Export Google Doc as DOCX
                file_content = drive_service.files().export_media(
                    fileId=google_drive_file_id,
                    mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                ).execute()
                file_ext = 'docx'
                if not filename.endswith('.docx'):
                    filename += '.docx'
            else:
                file_content = drive_service.files().get_media(fileId=google_drive_file_id).execute()
                file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

            if file_ext not in {'pdf', 'docx', 'doc', 'txt'}:
                return jsonify({'error': f'Unsupported file type: .{file_ext}'}), 400

            # Write to temp file for text extraction
            temp_path = os.path.join('/tmp', f"drive_{google_drive_file_id}.{file_ext}")
            with open(temp_path, 'wb') as f:
                f.write(file_content)

            extracted_text = ''
            try:
                if file_ext == 'pdf':
                    from pypdf import PdfReader
                    reader = PdfReader(temp_path)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text += page_text + '\n'
                elif file_ext in ['docx', 'doc']:
                    from docx import Document
                    doc = Document(temp_path)
                    for para in doc.paragraphs:
                        if para.text.strip():
                            extracted_text += para.text + '\n'
                elif file_ext == 'txt':
                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                        extracted_text = f.read()
            finally:
                if os.path.exists(temp_path):
                    os.remove(temp_path)

            if not extracted_text.strip():
                return jsonify({'error': f'Could not extract text from {filename}. File may be empty or image-based.'}), 400

            supabase = get_supabase_client()

            # Upsert: remove existing Drive record for same file_id, then insert
            supabase.table('documents').delete().eq('client_id', client_id).eq('drive_file_id', google_drive_file_id).execute()
            insert_result = supabase.table('documents').insert({
                'client_id': client_id,
                'filename': filename,
                'file_type': file_ext,
                'drive_file_id': google_drive_file_id,
                'source': 'google_drive',
                'extracted_text': extracted_text,
                'file_size': file_size or len(file_content),
                'drive_modified_time': drive_modified_time
            }).execute()

            # Deduct 2 credits via RPC
            supabase.rpc('deduct_credits', {
                'p_user_id': client_id,
                'p_amount': 2,
                'p_operation': 'google_drive_import',
                'p_reference_id': None
            }).execute()

            # RAG chunking — run in background so HTTP response is not delayed
            doc_id = insert_result.data[0]['id']
            threading.Thread(
                target=_embed_in_background,
                args=(doc_id, client_id, filename, extracted_text),
                daemon=True
            ).start()

            return jsonify({
                'success': True,
                'documents_processed': 1,
                'documents': [{'filename': filename, 'file_type': file_ext, 'chars_extracted': len(extracted_text)}]
            })

        # ── Manual file upload path (existing code) ───────────────────────────
        extract_only = request.form.get('extract_only') == 'true'
        files = request.files.getlist('files')

        if not files or all(f.filename == '' for f in files):
            return jsonify({'error': 'No files provided'}), 400

        if len(files) > 5:
            return jsonify({'error': 'Maximum 5 files allowed'}), 400

        MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
        ALLOWED_EXTENSIONS = {'pdf', 'docx', 'doc', 'txt'}

        supabase = get_supabase_client()
        processed = []

        for file in files:
            filename = file.filename
            if not filename:
                continue

            file_ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

            if file_ext not in ALLOWED_EXTENSIONS:
                return jsonify({'error': f'Unsupported file type: .{file_ext}. Allowed: PDF, DOCX, TXT'}), 400

            # Save to temp location
            temp_path = os.path.join('/tmp', f"upload_{datetime.now().strftime('%Y%m%d%H%M%S')}_{filename}")
            file.save(temp_path)
            file_size = os.path.getsize(temp_path)

            if file_size > MAX_FILE_SIZE:
                os.remove(temp_path)
                return jsonify({'error': f'{filename} exceeds 10MB limit'}), 400

            # Extract text based on file type
            extracted_text = ''
            try:
                if file_ext == 'pdf':
                    from pypdf import PdfReader
                    reader = PdfReader(temp_path)
                    for page in reader.pages:
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text += page_text + '\n'

                elif file_ext in ['docx', 'doc']:
                    from docx import Document
                    doc = Document(temp_path)
                    for para in doc.paragraphs:
                        if para.text.strip():
                            extracted_text += para.text + '\n'

                elif file_ext == 'txt':
                    with open(temp_path, 'r', encoding='utf-8', errors='ignore') as f:
                        extracted_text = f.read()

            finally:
                os.remove(temp_path)

            if not extracted_text.strip():
                return jsonify({'error': f'Could not extract text from {filename}. File may be empty or image-based.'}), 400

            # extract_only mode: return text without saving to DB
            if extract_only:
                return jsonify({
                    'success': True,
                    'documents_processed': 1,
                    'documents': [{'filename': filename, 'text': extracted_text}]
                })

            # Delete existing record for same filename + client, then insert fresh
            drive_file_id = f"upload_{filename}"
            supabase.table('documents').delete().eq('client_id', client_id).eq('drive_file_id', drive_file_id).execute()
            insert_result = supabase.table('documents').insert({
                'client_id': client_id,
                'filename': filename,
                'file_type': file_ext,
                'drive_file_id': drive_file_id,
                'extracted_text': extracted_text,
                'file_size': file_size
            }).execute()

            # ── RAG: chunk + embed the document ──────────────────────────────
            try:
                from embeddings import embed_texts, chunk_text
                doc_id = insert_result.data[0]['id']

                # Remove any stale chunks for this document
                supabase.table('document_chunks').delete().eq('document_id', doc_id).execute()

                # Split into overlapping passages and generate embeddings
                chunks = chunk_text(extracted_text)
                embeddings = embed_texts(chunks)

                chunk_rows = [
                    {
                        'document_id': doc_id,
                        'client_id': client_id,
                        'chunk_index': i,
                        'chunk_text': c,
                        'embedding': e,
                    }
                    for i, (c, e) in enumerate(zip(chunks, embeddings))
                ]

                # Insert in batches of 100 to stay within request size limits
                for batch_start in range(0, len(chunk_rows), 100):
                    supabase.table('document_chunks').insert(
                        chunk_rows[batch_start:batch_start + 100]
                    ).execute()

                print(f"[RAG] {len(chunks)} chunks indexed for {filename}")
            except Exception as embed_err:
                # Non-fatal: Q&A falls back to full-text if chunks are missing
                print(f"[RAG] Warning: could not index chunks for {filename}: {embed_err}")

            processed.append({
                'filename': filename,
                'file_type': file_ext,
                'file_size': file_size,
                'chars_extracted': len(extracted_text),
                'preview': extracted_text.strip()[:200]
            })

        return jsonify({
            'success': True,
            'documents_processed': len(processed),
            'documents': processed
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- NotebookLM content generation (podcast / infographic) ---

def _notebooklm_job_worker(job_id, storyboard_id, content_type, options=None):
    """Background worker: generate NotebookLM content and update job status."""
    try:
        from generate_notebooklm import generate_podcast, generate_infographic, generate_video, cleanup_notebook
        supabase = get_supabase_client()

        # Update job → processing
        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Fetch storyboard + course data
        sb_result = supabase.table('storyboards').select('*, courses(*)').eq('id', storyboard_id).single().execute()
        storyboard = sb_result.data
        if not storyboard:
            raise Exception('Storyboard not found')

        course_data = storyboard['content_json']
        course_id = storyboard['course_id']
        client_id = storyboard['courses']['client_id']

        # Gather source document text
        docs_result = supabase.table('documents').select('extracted_text').eq('client_id', client_id).limit(10).execute()
        source_text = '\n\n'.join(
            doc['extracted_text'][:5000] for doc in (docs_result.data or []) if doc.get('extracted_text')
        )
        # Fall back to storyboard content if no docs
        if not source_text.strip():
            slides_text = '\n'.join(
                f"{s.get('title', '')}: {' '.join(s.get('bullets', []))}"
                for s in course_data.get('slides', [])
            )
            source_text = f"{course_data.get('title', '')}\n\n{slides_text}"

        # Look up an existing NotebookLM notebook for this course to reuse
        existing_notebook_id = None
        prev = supabase.table('generation_jobs') \
            .select('notebooklm_notebook_id') \
            .eq('course_id', course_id) \
            .not_.is_('notebooklm_notebook_id', 'null') \
            .order('created_at', desc=True) \
            .limit(1) \
            .execute()
        if prev.data:
            existing_notebook_id = prev.data[0]['notebooklm_notebook_id']
            print(f"[NotebookLM] Found existing notebook for course {course_id}: {existing_notebook_id}")

        # Generate the artifact
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        notebook_id = None

        if content_type == 'podcast':
            filename = f"podcast_{storyboard_id[:8]}_{timestamp}.mp3"
            filepath = os.path.join(OUTPUT_DIR, filename)
            notebook_id = generate_podcast(source_text, course_data, filepath, options=options, existing_notebook_id=existing_notebook_id)
            storage_path = f"podcasts/{filename}"
            content_type_header = 'audio/mpeg'
            file_type = 'podcast'
        elif content_type == 'video':
            filename = f"video_{storyboard_id[:8]}_{timestamp}.mp4"
            filepath = os.path.join(OUTPUT_DIR, filename)
            notebook_id = generate_video(source_text, course_data, filepath, options=options, existing_notebook_id=existing_notebook_id)
            storage_path = f"videos/{filename}"
            content_type_header = 'video/mp4'
            file_type = 'video'
        else:
            filename = f"infographic_{storyboard_id[:8]}_{timestamp}.png"
            filepath = os.path.join(OUTPUT_DIR, filename)
            print(f"[NotebookLM] Starting infographic generation with options={options}")
            notebook_id = generate_infographic(source_text, course_data, filepath, options=options, existing_notebook_id=existing_notebook_id)
            storage_path = f"infographics/{filename}"
            content_type_header = 'image/png'
            file_type = 'infographic'

        # Update notebook ID for tracking
        supabase.table('generation_jobs').update({
            'notebooklm_notebook_id': notebook_id,
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Upload to Supabase Storage
        with open(filepath, 'rb') as f:
            file_bytes = f.read()

        supabase.storage.from_('course-files').upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": content_type_header}
        )

        public_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Save to generated_files table
        file_record = supabase.table('generated_files').insert({
            'course_id': course_id,
            'file_type': file_type,
            'file_url': public_url,
            'file_size': len(file_bytes)
        }).execute()

        # Update job → completed
        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Cleanup local temp file
        os.remove(filepath)
        # NOTE: intentionally not deleting the NotebookLM notebook so the
        # notebook and its studio assets (audio, infographic, video) remain
        # visible in NotebookLM for inspection.
        # if notebook_id:
        #     try:
        #         cleanup_notebook(notebook_id)
        #     except Exception as e:
        #         print(f"[NotebookLM] Warning: cleanup failed for notebook {notebook_id}: {e}")

        print(f"[NotebookLM] {file_type} generated successfully: {filename}")

    except Exception as err:
        import traceback
        print(f"[NotebookLM] Job {job_id} failed: {err}")
        print(f"[NotebookLM] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
        except Exception:
            pass


@app.route('/generate-notebooklm-content', methods=['POST'])
def generate_notebooklm_content():
    """
    Start async generation of a podcast or infographic via NotebookLM.
    Expected JSON: { "storyboard_id": "uuid", "content_type": "podcast"|"infographic" }
    Returns: { "job_id": "uuid" } immediately; poll /job-status/<job_id> for result.
    """
    try:
        data = request.json or {}
        storyboard_id = data.get('storyboard_id')
        content_type = data.get('content_type')
        options = data.get('options', {})

        if not storyboard_id:
            return jsonify({'error': 'storyboard_id is required'}), 400
        if content_type not in ('podcast', 'infographic', 'video'):
            return jsonify({'error': 'content_type must be "podcast", "infographic", or "video"'}), 400

        # Look up course_id from the storyboard
        supabase = get_supabase_client()
        sb = supabase.table('storyboards').select('course_id').eq('id', storyboard_id).single().execute()
        if not sb.data:
            return jsonify({'error': 'Storyboard not found'}), 404

        # Create job row
        job = supabase.table('generation_jobs').insert({
            'course_id': sb.data['course_id'],
            'job_type': content_type,
            'status': 'pending'
        }).execute()

        job_id = job.data[0]['id']

        # Spawn background worker
        threading.Thread(
            target=_notebooklm_job_worker,
            args=(job_id, storyboard_id, content_type, options),
            daemon=True
        ).start()

        return jsonify({'success': True, 'job_id': job_id})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


def _slide_deck_job_worker(job_id, document_ids, user_id, options=None, existing_notebook_id=None):
    """Background worker: generate slide deck via NotebookLM or Claude, convert to preview images, upload all."""
    try:
        import shutil
        import time as _time
        from generate_notebooklm import generate_slide_deck
        job_start = _time.time()
        supabase = get_supabase_client()

        # Update job → processing
        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Fetch document texts by ID
        docs_result = supabase.table('documents') \
            .select('id, filename, extracted_text') \
            .in_('id', document_ids) \
            .execute()

        documents = [doc for doc in (docs_result.data or []) if doc.get('extracted_text', '').strip()]
        if not documents:
            raise Exception('No text content found in selected documents')

        title = (options or {}).get('title') or (documents[0].get('filename', 'Slide Deck'))

        # Create temp dir for this job's outputs
        job_output_dir = os.path.join(OUTPUT_DIR, f'slides_{job_id}')
        os.makedirs(job_output_dir, exist_ok=True)

        # Callback: persist notebook_id as soon as it's known so a retry is possible even if we time out
        def _save_notebook_id(nb_id):
            try:
                supabase.table('generation_jobs').update({
                    'notebooklm_notebook_id': nb_id,
                    'updated_at': datetime.now().isoformat()
                }).eq('id', job_id).execute()
                print(f"[NotebookLM] Saved notebook_id={nb_id} for job {job_id}")
            except Exception as e:
                print(f"[NotebookLM] Failed to save notebook_id (non-fatal): {e}")

        # Generate via selected engine, with automatic fallback
        engine = options.get('slide_engine', 'notebooklm')
        used_fallback = False

        if engine == 'claude':
            from generate_claude_slides import generate_slide_deck_claude
            print(f"[Slides] Using Claude engine for job {job_id}")
            notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck_claude(
                documents, title, job_output_dir, options=options
            )
        else:
            print(f"[Slides] Using NotebookLM engine for job {job_id}")
            try:
                notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck(
                    documents, title, job_output_dir,
                    options=options, existing_notebook_id=existing_notebook_id,
                    notebook_id_callback=_save_notebook_id
                )
            except Exception as nblm_err:
                print(f"[Slides] NotebookLM failed: {nblm_err} — falling back to Claude")
                from generate_claude_slides import generate_slide_deck_claude
                used_fallback = True
                engine = 'claude'
                notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck_claude(
                    documents, title, job_output_dir, options=options
                )

        # Post-process PPTX: white strip + speaker notes from voiceover scripts
        if pptx_path and os.path.exists(pptx_path):
            try:
                from pptx import Presentation
                from pptx.util import Inches, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
                prs = Presentation(pptx_path)
                strip_height = Inches(0.45)
                for i, slide in enumerate(prs.slides):
                    # Branding strip
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                        left=Emu(0),
                        top=Emu(prs.slide_height - int(strip_height)),
                        width=Emu(prs.slide_width),
                        height=strip_height,
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    shape.line.fill.background()
                    # Speaker notes from voiceover script
                    if voiceover_scripts and i < len(voiceover_scripts):
                        script = (voiceover_scripts[i] or '').strip()
                        if script:
                            notes_slide = slide.notes_slide
                            tf = notes_slide.notes_text_frame
                            tf.text = script
                prs.save(pptx_path)
                print(f"[PPTX] Branding strip + speaker notes applied to {len(prs.slides)} slides")
            except Exception as e:
                print(f"[PPTX] Post-processing failed ({e}), uploading original")

        # Update notebook ID for tracking
        supabase.table('generation_jobs').update({
            'notebooklm_notebook_id': notebook_id,
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')

        # Upload PDF to Supabase Storage (may not exist for Claude engine)
        pdf_url = None
        if pdf_path and os.path.exists(pdf_path):
            pdf_storage_path = f"slide_decks/{job_id}_{timestamp}.pdf"
            with open(pdf_path, 'rb') as f:
                supabase.storage.from_('course-files').upload(
                    path=pdf_storage_path, file=f.read(),
                    file_options={"content-type": "application/pdf"}
                )
            pdf_url = supabase.storage.from_('course-files').get_public_url(pdf_storage_path)

        # Upload PPTX if available
        pptx_url = None
        pptx_size = 0
        if pptx_path and os.path.exists(pptx_path):
            pptx_storage_path = f"slide_decks/{job_id}_{timestamp}.pptx"
            with open(pptx_path, 'rb') as f:
                pptx_bytes = f.read()
                pptx_size = len(pptx_bytes)
                supabase.storage.from_('course-files').upload(
                    path=pptx_storage_path, file=pptx_bytes,
                    file_options={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
                )
            pptx_url = supabase.storage.from_('course-files').get_public_url(pptx_storage_path)

        # Upload per-slide preview images
        slide_image_urls = []
        for i, img_path in enumerate(slide_image_paths):
            img_storage_path = f"slide_decks/{job_id}_slide_{i+1}.png"
            with open(img_path, 'rb') as f:
                supabase.storage.from_('course-files').upload(
                    path=img_storage_path, file=f.read(),
                    file_options={"content-type": "image/png"}
                )
            url = supabase.storage.from_('course-files').get_public_url(img_storage_path)
            slide_image_urls.append(url)

        # Use PPTX URL as primary file_url, fall back to PDF
        primary_url = pptx_url or pdf_url
        primary_size = pptx_size or (os.path.getsize(pdf_path) if pdf_path and os.path.exists(pdf_path) else 0)

        # Save to generated_files table
        elapsed_secs = round(_time.time() - job_start, 1)
        file_record = supabase.table('generated_files').insert({
            'file_type': 'slide_deck',
            'file_url': primary_url,
            'file_size': primary_size,
            'metadata': json.dumps({
                'slide_images': slide_image_urls,
                'pdf_url': pdf_url,
                'pptx_url': pptx_url,
                'notebook_id': notebook_id,
                'slide_count': len(slide_image_urls),
                'voiceover_scripts': voiceover_scripts,
                'target_time': options.get('target_time') if options else None,
                'max_time': options.get('max_time') if options else None,
                'title': title,
                'slide_engine': engine,
                'used_fallback': used_fallback,
                'generation_time_secs': elapsed_secs,
            })
        }).execute()

        # Update job → completed
        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Cleanup local temp files
        shutil.rmtree(job_output_dir, ignore_errors=True)

        engine = options.get('slide_engine', 'notebooklm') if options else 'notebooklm'
        print(f"[Slides] Deck generated via {engine}: {len(slide_image_urls)} slides in {elapsed_secs}s")

    except Exception as err:
        import traceback
        print(f"[NotebookLM] Slide deck job {job_id} failed: {err}")
        print(f"[NotebookLM] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
        except Exception:
            pass


@app.route('/auto-complete-course', methods=['POST'])
def auto_complete_course():
    """
    Start async auto-complete pipeline: slides → quiz → audio → video → (scorm).
    Expected JSON: { "document_ids": [...], "user_id": "uuid", "target": "mp4"|"scorm",
                     "credit_cost": int, "options": {...} }
    Returns: { "job_id": "uuid" } immediately; poll /job-status/<job_id>.
    """
    try:
        data = request.json or {}
        document_ids = data.get('document_ids', [])
        user_id = data.get('user_id')
        target = data.get('target', 'mp4')
        credit_cost = data.get('credit_cost', 8)
        options = data.get('options', {})

        if not document_ids:
            return jsonify({'error': 'document_ids is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        supabase = get_supabase_client()
        job = supabase.table('generation_jobs').insert({
            'job_type': 'auto_complete',
            'status': 'pending',
            'user_id': user_id,
        }).execute()
        job_id = job.data[0]['id']

        threading.Thread(
            target=_auto_complete_worker,
            args=(job_id, document_ids, user_id, target, credit_cost, options),
            daemon=True
        ).start()

        return jsonify({'success': True, 'job_id': job_id})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


def _auto_complete_worker(job_id, document_ids, user_id, target, credit_cost, options):
    """Background worker: runs full pipeline slides → quiz → audio → video → (scorm)."""
    try:
        import shutil
        import time as _time
        import requests as http_requests
        from generate_notebooklm import generate_slide_deck, _extract_slide_texts, _clean_voiceover_script, _parse_voiceover_scripts

        job_start = _time.time()
        supabase = get_supabase_client()
        completed_steps = []
        meta = {'target': target, 'current_step': 'slides', 'steps_completed': [], 'steps_remaining': ['slides', 'quiz', 'audio', 'video'] + (['scorm'] if target == 'scorm' else [])}

        def _update_progress(step, extra_meta=None):
            meta['current_step'] = step
            remaining = [s for s in ['slides', 'quiz', 'audio', 'video', 'scorm'] if s not in completed_steps and s != step]
            if target != 'scorm':
                remaining = [s for s in remaining if s != 'scorm']
            meta['steps_completed'] = list(completed_steps)
            meta['steps_remaining'] = remaining
            if extra_meta:
                meta.update(extra_meta)
            supabase.table('generation_jobs').update({
                'status': 'processing',
                'updated_at': datetime.now().isoformat(),
            }).eq('id', job_id).execute()
            # Update metadata in a temp generated_files row if exists
            print(f"[auto-complete] Step: {step} | Completed: {completed_steps}")

        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # ── Step 1: Generate slides ──
        _update_progress('slides')
        docs_result = supabase.table('documents') \
            .select('id, filename, extracted_text') \
            .in_('id', document_ids) \
            .execute()
        documents = [d for d in (docs_result.data or []) if d.get('extracted_text')]
        if not documents:
            raise Exception('No documents with text found')

        title = options.get('title', 'Training Course')
        job_output_dir = os.path.join(OUTPUT_DIR, f'auto_{job_id}')
        os.makedirs(job_output_dir, exist_ok=True)

        engine = options.get('slide_engine', 'notebooklm')
        if engine == 'claude':
            from generate_claude_slides import generate_slide_deck_claude
            notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck_claude(
                documents, title, job_output_dir, options=options
            )
        else:
            try:
                notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck(
                    documents, title, job_output_dir, options=options
                )
            except Exception:
                from generate_claude_slides import generate_slide_deck_claude
                notebook_id, pdf_path, pptx_path, slide_image_paths, voiceover_scripts = generate_slide_deck_claude(
                    documents, title, job_output_dir, options=options
                )

        # Upload slides
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        pptx_url = None
        if pptx_path and os.path.exists(pptx_path):
            pptx_storage = f"slide_decks/{job_id}_{timestamp}.pptx"
            with open(pptx_path, 'rb') as f:
                supabase.storage.from_('course-files').upload(path=pptx_storage, file=f.read(),
                    file_options={"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"})
            pptx_url = supabase.storage.from_('course-files').get_public_url(pptx_storage)

        slide_image_urls = []
        for i, img_path in enumerate(slide_image_paths or []):
            img_storage = f"slide_decks/{job_id}_slide_{i+1}.png"
            with open(img_path, 'rb') as f:
                supabase.storage.from_('course-files').upload(path=img_storage, file=f.read(),
                    file_options={"content-type": "image/png"})
            slide_image_urls.append(supabase.storage.from_('course-files').get_public_url(img_storage))

        completed_steps.append('slides')

        # ── Step 2: Generate quiz ──
        _update_progress('quiz')
        quiz_data = None
        try:
            anthropic_key = os.environ.get('ANTHROPIC_API_KEY')
            if anthropic_key and voiceover_scripts:
                slide_content = '\n\n'.join(f'[Slide {i+1}] {s}' for i, s in enumerate(voiceover_scripts) if s)
                doc_context = '\n\n---\n\n'.join(f'Document: {d["filename"]}\n{d.get("extracted_text", "")[:30000]}' for d in documents[:3])
                count = options.get('quiz_count', 10)
                difficulty = options.get('quiz_difficulty', 'medium')
                pass_pct = options.get('quiz_pass_percent', 70)

                quiz_prompt = f"""Generate exactly {count} multiple-choice quiz questions based on this training content.
Difficulty: {difficulty}. Each question has 4 options (A-D), one correct answer.

--- SLIDES ---
{slide_content}

--- SOURCE DOCS ---
{doc_context[:50000]}

Return ONLY valid JSON:
{{"questions": [{{"question": "...", "options": ["A...", "B...", "C...", "D..."], "correct": 0, "explanation": "..."}}], "pass_percentage": {pass_pct}}}"""

                resp = http_requests.post('https://api.anthropic.com/v1/messages',
                    headers={'Content-Type': 'application/json', 'x-api-key': anthropic_key, 'anthropic-version': '2023-06-01'},
                    json={'model': 'claude-sonnet-4-6', 'max_tokens': 8000, 'messages': [{'role': 'user', 'content': quiz_prompt}]},
                    timeout=120)
                if resp.ok:
                    txt = resp.json()['content'][0]['text']
                    js = txt[txt.index('{'):txt.rindex('}')+1]
                    quiz_data = json.loads(js)
                    print(f"[auto-complete] Quiz: {len(quiz_data.get('questions', []))} questions generated")
        except Exception as qe:
            print(f"[auto-complete] Quiz generation failed (non-fatal): {qe}")
        completed_steps.append('quiz')

        # ── Step 3: Generate audio ──
        _update_progress('audio')
        voiceover_zip_url = None
        if voiceover_scripts and any(s.strip() for s in voiceover_scripts):
            voice_id = options.get('voice_id', 'aura-athena-en')
            voice_provider = options.get('voice_provider', 'deepgram')

            from pydub import AudioSegment
            import zipfile
            import io
            from concurrent.futures import ThreadPoolExecutor, as_completed

            audio_dir = os.path.join(job_output_dir, 'audio')
            os.makedirs(audio_dir, exist_ok=True)

            deepgram_key = os.environ.get('DEEPGRAM_API_KEY')
            alibaba_key = os.environ.get('ALIBABA_API_KEY')

            def _tts_chunk(text_chunk):
                if voice_provider == 'deepgram' and deepgram_key:
                    r = http_requests.post(f'https://api.deepgram.com/v1/speak?model={voice_id}',
                        headers={'Authorization': f'Token {deepgram_key}', 'Content-Type': 'application/json'},
                        json={"text": text_chunk}, timeout=120)
                    r.raise_for_status()
                    return ('mp3', r.content)
                else:
                    r = http_requests.post('https://dashscope-intl.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                        headers={'Authorization': f'Bearer {alibaba_key}', 'Content-Type': 'application/json'},
                        json={"model": "qwen3-tts-vc-2026-01-22", "input": {"text": text_chunk, "voice": voice_id}}, timeout=120)
                    r.raise_for_status()
                    audio_url = r.json().get('output', {}).get('audio', {}).get('url')
                    wav = http_requests.get(audio_url, timeout=60)
                    return ('wav', wav.content)

            mp3_paths = []
            for si, script in enumerate(voiceover_scripts):
                script = (script or '').strip()
                if not script:
                    silence = AudioSegment.silent(duration=1000)
                    mp3p = os.path.join(audio_dir, f'slide_{si+1:02d}.mp3')
                    silence.export(mp3p, format='mp3')
                    mp3_paths.append(mp3p)
                    continue

                from generate_notebooklm import chunk_text_for_tts
                chunks = chunk_text_for_tts(script, max_chars=512)
                chunk_results = [None] * len(chunks)
                with ThreadPoolExecutor(max_workers=3) as executor:
                    futs = {executor.submit(_tts_chunk, c): i for i, c in enumerate(chunks)}
                    for fut in as_completed(futs):
                        chunk_results[futs[fut]] = fut.result()

                combined = AudioSegment.empty()
                for fmt, audio_bytes in chunk_results:
                    seg = AudioSegment.from_mp3(io.BytesIO(audio_bytes)) if fmt == 'mp3' else AudioSegment.from_wav(io.BytesIO(audio_bytes))
                    combined += seg

                mp3p = os.path.join(audio_dir, f'slide_{si+1:02d}.mp3')
                combined.export(mp3p, format='mp3', bitrate='128k')
                mp3_paths.append(mp3p)

            # Zip audio
            zip_filename = f'voiceover_{job_id[:8]}_{timestamp}.zip'
            zip_path = os.path.join(job_output_dir, zip_filename)
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for mp3p in mp3_paths:
                    zf.write(mp3p, os.path.basename(mp3p))

            vo_storage = f'voiceovers/{zip_filename}'
            with open(zip_path, 'rb') as f:
                supabase.storage.from_('course-files').upload(path=vo_storage, file=f.read(),
                    file_options={"content-type": "application/zip"})
            voiceover_zip_url = supabase.storage.from_('course-files').get_public_url(vo_storage)
            print(f"[auto-complete] Audio: {len(mp3_paths)} slide MP3s zipped")

        completed_steps.append('audio')

        # ── Step 4: Generate video ──
        _update_progress('video')
        video_url = None
        if voiceover_zip_url and slide_image_urls:
            # Call existing video endpoint internally
            video_job = supabase.table('generation_jobs').insert({
                'job_type': 'voiceover_video', 'status': 'pending', 'user_id': user_id,
            }).execute()
            vid_job_id = video_job.data[0]['id']

            _voiceover_video_job_worker(vid_job_id, voiceover_zip_url, slide_image_urls, user_id)

            # Read result
            vid_result = supabase.table('generation_jobs').select('status, result_file_id').eq('id', vid_job_id).single().execute()
            if vid_result.data and vid_result.data['status'] == 'completed' and vid_result.data.get('result_file_id'):
                file_data = supabase.table('generated_files').select('file_url').eq('id', vid_result.data['result_file_id']).single().execute()
                if file_data.data:
                    video_url = file_data.data['file_url']
                    print(f"[auto-complete] Video: {video_url}")
            else:
                print(f"[auto-complete] Video generation failed — continuing")

        completed_steps.append('video')

        # ── Step 5: Generate SCORM (if target=scorm) ──
        scorm_url = None
        if target == 'scorm' and video_url:
            _update_progress('scorm')
            try:
                from generate_scorm import generate_scorm_package, generate_api_wrapper

                scorm_dir = os.path.join(job_output_dir, 'scorm')
                os.makedirs(scorm_dir, exist_ok=True)

                course_data = {'title': title, 'slides': []}
                scorm_zip_path = os.path.join(scorm_dir, 'scorm_package.zip')

                # Build minimal course data for SCORM with video embed
                import zipfile
                with zipfile.ZipFile(scorm_zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                    # imsmanifest.xml
                    manifest = f'''<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="com.trusource.course" version="1.0"
  xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
  xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2">
  <organizations default="org1">
    <organization identifier="org1"><title>{title}</title>
      <item identifier="item1" identifierref="res1"><title>{title}</title></item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="res1" type="webcontent" adlcp:scormtype="sco" href="index.html">
      <file href="index.html"/>
      <file href="APIWrapper.js"/>
    </resource>
  </resources>
</manifest>'''
                    zf.writestr('imsmanifest.xml', manifest)
                    zf.writestr('APIWrapper.js', generate_api_wrapper())

                    # Build HTML with video + quiz
                    quiz_html = ''
                    if quiz_data and quiz_data.get('questions'):
                        pass_pct = quiz_data.get('pass_percentage', 70)
                        quiz_json_str = json.dumps(quiz_data['questions'])
                        quiz_html = f'''
<div style="margin-top:30px;padding:20px;background:#f3e5f5;border:2px solid #ce93d8;border-radius:12px;">
<h2 style="color:#7b1fa2;">Knowledge Check</h2>
<p style="color:#666;">Answer all questions, then click Submit. Pass mark: {pass_pct}%</p>
<div id="quiz-container"></div>
<div id="quiz-result" style="display:none;padding:20px;border-radius:12px;margin-top:20px;text-align:center;"></div>
<button id="quiz-submit" onclick="submitQuiz()" style="margin-top:15px;background:#7b1fa2;color:white;border:none;padding:12px 32px;border-radius:8px;font-size:16px;font-weight:600;cursor:pointer;">Submit Answers</button>
</div>
<script>
(function(){{var questions={quiz_json_str};var passPercent={pass_pct};var c=document.getElementById("quiz-container");
questions.forEach(function(q,i){{var d=document.createElement("div");d.style.cssText="margin-bottom:20px;padding:16px;background:white;border-radius:10px;border:1px solid #e0e0e0;";
var h='<p style="font-weight:600;margin-bottom:10px;">'+(i+1)+". "+q.question+"</p>";
q.options.forEach(function(o,j){{var l=String.fromCharCode(65+j);h+='<label style="display:block;padding:8px 12px;margin:4px 0;border-radius:6px;cursor:pointer;border:1px solid #e0e0e0;"><input type="radio" name="q'+i+'" value="'+j+'" style="margin-right:8px;"> '+l+". "+o+"</label>";}});
d.innerHTML=h;c.appendChild(d);}});
window.submitQuiz=function(){{var score=0;var total=questions.length;
questions.forEach(function(q,i){{var s=document.querySelector('input[name="q'+i+'"]:checked');if(s&&parseInt(s.value)===q.correct)score++;}});
var pct=Math.round((score/total)*100);var passed=pct>=passPercent;var r=document.getElementById("quiz-result");r.style.display="block";
r.style.background=passed?"#e8f5e9":"#ffebee";r.style.border="2px solid "+(passed?"#4caf50":"#f44336");
r.innerHTML='<h2 style="color:'+(passed?"#2e7d32":"#c62828")+';margin:0 0 8px 0;">'+(passed?"PASSED":"NOT PASSED")+"</h2><p style='font-size:24px;font-weight:bold;margin:0;'>"+score+" / "+total+" ("+pct+"%)</p><p style='color:#666;margin:8px 0 0 0;'>Pass mark: "+passPercent+"%</p>";
document.getElementById("quiz-submit").style.display="none";
if(window.API){{API.LMSSetValue("cmi.core.score.raw",String(pct));API.LMSSetValue("cmi.core.lesson_status",passed?"passed":"failed");API.LMSCommit("");}}}};
}})();
</script>'''

                    index_html = f'''<!DOCTYPE html><html><head><meta charset="utf-8"><title>{title}</title>
<script src="APIWrapper.js"></script>
<style>body{{font-family:-apple-system,sans-serif;max-width:900px;margin:0 auto;padding:20px;background:#fafafa;}}
h1{{color:#1a1a2e;}}video{{width:100%;border-radius:12px;margin:20px 0;}}</style></head>
<body><h1>{title}</h1>
<video controls><source src="{video_url}" type="video/mp4">Your browser does not support video.</video>
{quiz_html}
<button onclick="if(window.API){{API.LMSSetValue('cmi.core.lesson_status','completed');API.LMSCommit('');alert('Course completed!');}}"
style="margin-top:20px;background:#0a84ff;color:white;border:none;padding:12px 32px;border-radius:8px;font-size:16px;font-weight:600;cursor:pointer;">Complete Course</button>
<script>if(window.API){{API.LMSInitialize("");API.LMSSetValue("cmi.core.lesson_status","incomplete");}}</script>
</body></html>'''
                    zf.writestr('index.html', index_html)

                scorm_storage = f'scorm/{job_id}_{timestamp}.zip'
                with open(scorm_zip_path, 'rb') as f:
                    supabase.storage.from_('course-files').upload(path=scorm_storage, file=f.read(),
                        file_options={"content-type": "application/zip"})
                scorm_url = supabase.storage.from_('course-files').get_public_url(scorm_storage)
                print(f"[auto-complete] SCORM: {scorm_url}")
                completed_steps.append('scorm')
            except Exception as se:
                print(f"[auto-complete] SCORM generation failed: {se}")

        # ── Save final result ──
        elapsed_secs = round(_time.time() - job_start, 1)
        file_record = supabase.table('generated_files').insert({
            'file_type': 'auto_complete',
            'file_url': video_url or pptx_url or '',
            'file_size': 0,
            'metadata': json.dumps({
                'target': target,
                'slide_images': slide_image_urls,
                'pptx_url': pptx_url,
                'voiceover_scripts': voiceover_scripts,
                'voiceover_zip_url': voiceover_zip_url,
                'video_url': video_url,
                'scorm_url': scorm_url,
                'quiz': quiz_data,
                'title': title,
                'slide_engine': engine,
                'generation_time_secs': elapsed_secs,
                'steps_completed': completed_steps,
            })
        }).execute()

        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        shutil.rmtree(job_output_dir, ignore_errors=True)
        print(f"[auto-complete] Pipeline complete: {completed_steps} in {elapsed_secs}s")

    except Exception as err:
        import traceback
        print(f"[auto-complete] Job {job_id} failed: {err}")
        print(f"[auto-complete] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
            # Refund credits for uncompleted steps
            step_costs = {'slides': 3, 'quiz': 2, 'audio': 3, 'video': 5 if target == 'scorm' else 5, 'scorm': 2}
            refund = sum(v for k, v in step_costs.items() if k not in completed_steps and (k != 'scorm' or target == 'scorm'))
            if refund > 0:
                supabase.rpc('refund_credits', {'p_user_id': user_id, 'p_amount': refund, 'p_operation': 'auto_complete_refund', 'p_reference_id': None}).execute()
                print(f"[auto-complete] Refunded {refund} credits for failed steps")
        except Exception:
            pass


@app.route('/generate-slides-content', methods=['POST'])
def generate_slides_content():
    """
    Start async slide deck generation via NotebookLM.
    Expected JSON: { "document_ids": [...], "user_id": "uuid",
                     "instructions": "...", "slide_format": "DETAILED_DECK"|"PRESENTER_SLIDES",
                     "slide_length": "DEFAULT"|"SHORT",
                     "existing_notebook_id": null|"..." }
    Returns: { "job_id": "uuid" } immediately; poll /job-status/<job_id>.
    """
    try:
        data = request.json or {}
        document_ids = data.get('document_ids', [])
        user_id = data.get('user_id')
        # Read options from nested object or flat params (backwards compat)
        opts = data.get('options') or {}
        options = {
            'slide_engine': opts.get('slide_engine', data.get('slide_engine', 'notebooklm')),
            'instructions': opts.get('instructions', data.get('instructions')),
            'slide_format': opts.get('slide_format', data.get('slide_format', 'DETAILED_DECK')),
            'slide_length': opts.get('slide_length', data.get('slide_length', 'DEFAULT')),
            'target_time': opts.get('target_time', data.get('target_time')),
            'max_time': opts.get('max_time', data.get('max_time')),
            'ppt_theme': opts.get('ppt_theme', data.get('ppt_theme')),
            'ppt_template_url': opts.get('ppt_template_url', data.get('ppt_template_url')),
            'title': opts.get('title', data.get('title')),
            'slide_count': opts.get('slide_count', data.get('slide_count')),
        }
        existing_notebook_id = opts.get('existing_notebook_id', data.get('existing_notebook_id'))

        if not document_ids:
            return jsonify({'error': 'document_ids is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        supabase = get_supabase_client()

        # Create job row (no course_id for standalone slide generation)
        job = supabase.table('generation_jobs').insert({
            'job_type': 'slide_deck',
            'status': 'pending',
            'user_id': user_id,
        }).execute()

        job_id = job.data[0]['id']

        # Spawn background worker
        threading.Thread(
            target=_slide_deck_job_worker,
            args=(job_id, document_ids, user_id, options, existing_notebook_id),
            daemon=True
        ).start()

        return jsonify({'success': True, 'job_id': job_id})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/job-status/<job_id>', methods=['GET'])
def job_status(job_id):
    """
    Poll the status of a NotebookLM generation job.
    Returns: { "status": "pending"|"processing"|"completed"|"failed",
               "file_url": "...", "error_message": "..." }
    """
    try:
        supabase = get_supabase_client()
        result = supabase.table('generation_jobs') \
            .select('status, error_message, result_file_id') \
            .eq('id', job_id) \
            .single() \
            .execute()

        if not result.data:
            return jsonify({'error': 'Job not found'}), 404

        job = result.data
        response = {'status': job['status']}

        if job['status'] == 'completed' and job.get('result_file_id'):
            file_result = supabase.table('generated_files') \
                .select('file_url, metadata') \
                .eq('id', job['result_file_id']) \
                .single() \
                .execute()
            if file_result.data:
                response['file_url'] = file_result.data['file_url']
                # Include metadata for slide deck jobs (contains slide_images, pdf_url, etc.)
                if file_result.data.get('metadata'):
                    meta = file_result.data['metadata']
                    if isinstance(meta, str):
                        try:
                            meta = json.loads(meta)
                        except Exception:
                            meta = None
                    if meta:
                        response['metadata'] = meta

        if job.get('error_message'):
            response['error_message'] = job['error_message']

        return jsonify(response)

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/notebooklm-status', methods=['GET'])
def notebooklm_status():
    """
    Health check for NotebookLM auth.
    Returns: { "available": true|false, "error": "..." }
    """
    auth_json = os.environ.get('NOTEBOOKLM_AUTH_JSON')
    if not auth_json:
        return jsonify({'available': False, 'error': 'NOTEBOOKLM_AUTH_JSON not configured'})

    try:
        from generate_notebooklm import check_auth
        check_auth()
        return jsonify({'available': True})
    except Exception as e:
        return jsonify({'available': False, 'error': str(e)})


# ── TTS Voice Cloning & Voiceover Audio ──────────────────────────────────────

def chunk_text_for_tts(text, max_chars=512):
    """Split text into ≤512-char chunks at sentence boundaries."""
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    chunks, current = [], ''
    for sentence in sentences:
        if not sentence.strip():
            continue
        # If single sentence exceeds limit, split at word boundaries
        if len(sentence) > max_chars:
            if current.strip():
                chunks.append(current.strip())
                current = ''
            words = sentence.split(' ')
            sub = ''
            for w in words:
                test = (sub + ' ' + w).strip()
                if len(test) <= max_chars:
                    sub = test
                else:
                    if sub.strip():
                        chunks.append(sub.strip())
                    sub = w
            if sub.strip():
                chunks.append(sub.strip())
            continue
        # Normal case: accumulate sentences
        test = (current + ' ' + sentence).strip() if current else sentence
        if len(test) <= max_chars:
            current = test
        else:
            if current.strip():
                chunks.append(current.strip())
            current = sentence
    if current.strip():
        chunks.append(current.strip())
    return chunks


@app.route('/enroll-voice', methods=['POST'])
def enroll_voice():
    """
    Enroll a voice with Qwen TTS voice cloning.
    Expected JSON: { "audio_base64": "data:audio/wav;base64,...",
                     "transcript": "...", "name": "My Voice", "user_id": "uuid" }
    Returns: { "success": true, "voice_id": "...", "name": "..." }
    """
    try:
        import requests as http_requests
        data = request.json or {}
        audio_base64 = data.get('audio_base64')
        transcript = data.get('transcript', '')
        name = data.get('name', 'My Voice')
        user_id = data.get('user_id')

        if not audio_base64:
            return jsonify({'error': 'audio_base64 is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        alibaba_key = os.environ.get('ALIBABA_API_KEY')
        if not alibaba_key:
            return jsonify({'error': 'ALIBABA_API_KEY not configured'}), 500

        enrollment_payload = {
            "model": "qwen-voice-enrollment",
            "input": {
                "action": "create",
                "target_model": "qwen3-tts-vc-2026-01-22",
                "preferred_name": f"voice_{user_id[:8]}",
                "audio": {
                    "data": audio_base64,
                },
            }
        }
        if transcript.strip():
            enrollment_payload["input"]["text"] = transcript.strip()

        resp = http_requests.post(
            'https://dashscope-intl.aliyuncs.com/api/v1/services/audio/tts/customization',
            headers={
                'Authorization': f'Bearer {alibaba_key}',
                'Content-Type': 'application/json',
            },
            json=enrollment_payload,
            timeout=60,
        )
        resp.raise_for_status()
        result = resp.json()
        voice_id = result.get('output', {}).get('voice_id') or result.get('output', {}).get('voice')

        if not voice_id:
            return jsonify({'error': 'Enrollment failed: no voice_id returned', 'raw': result}), 500

        # Save to tts_voices table
        supabase = get_supabase_client()
        voice_record = supabase.table('tts_voices').insert({
            'user_id': user_id,
            'voice_id': voice_id,
            'name': name,
        }).execute()

        print(f"[TTS] Voice enrolled: {name} -> {voice_id}")

        return jsonify({
            'success': True,
            'voice_id': voice_id,
            'record_id': voice_record.data[0]['id'],
            'name': name,
        })

    except Exception as e:
        import traceback
        print(f"[TTS] Voice enrollment failed: {e}")
        print(f"[TTS] Traceback:\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500


@app.route('/list-voices', methods=['POST'])
def list_voices():
    """List enrolled TTS voices for a user."""
    try:
        data = request.json or {}
        user_id = data.get('user_id')
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        supabase = get_supabase_client()
        result = supabase.table('tts_voices') \
            .select('id, voice_id, name, is_default, created_at') \
            .eq('user_id', user_id) \
            .order('created_at', desc=True) \
            .execute()

        return jsonify({'success': True, 'voices': result.data or []})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/preview-voice', methods=['GET', 'POST'])
def preview_voice():
    """Generate a short TTS preview. Supports Deepgram (stock) and Qwen (cloned) voices."""
    try:
        import requests as http_requests
        if request.method == 'GET':
            voice_id = request.args.get('voice_id', 'aura-asteria-en')
            provider = request.args.get('provider', 'deepgram')
            text = 'Welcome to your training course. This is a preview of the selected voice.'
        else:
            data = request.json or {}
            voice_id = data.get('voice_id', 'aura-asteria-en')
            provider = data.get('provider', 'deepgram')
            text = data.get('text', 'Welcome to your training course.')

        from flask import Response as FlaskResponse

        if provider == 'qwen':
            alibaba_key = os.environ.get('ALIBABA_API_KEY')
            if not alibaba_key:
                return jsonify({'error': 'ALIBABA_API_KEY not configured'}), 500

            resp = http_requests.post(
                'https://dashscope-intl.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                headers={'Authorization': f'Bearer {alibaba_key}', 'Content-Type': 'application/json'},
                json={"model": "qwen3-tts-vc-2026-01-22", "input": {"text": text, "voice": voice_id}},
                timeout=60,
            )
            resp.raise_for_status()
            result = resp.json()
            audio_url = result.get('output', {}).get('audio', {}).get('url')
            if not audio_url:
                return jsonify({'error': 'No audio URL returned'}), 500
            wav_resp = http_requests.get(audio_url, timeout=30)
            wav_resp.raise_for_status()
            # Convert WAV to MP3 for browser playback
            from pydub import AudioSegment
            import io
            audio = AudioSegment.from_wav(io.BytesIO(wav_resp.content))
            mp3_buf = io.BytesIO()
            audio.export(mp3_buf, format='mp3', bitrate='128k')
            return FlaskResponse(mp3_buf.getvalue(), mimetype='audio/mpeg', headers={
                'Access-Control-Allow-Origin': '*',
                'Cache-Control': 'public, max-age=3600',
            })

        # Default: Deepgram
        deepgram_key = os.environ.get('DEEPGRAM_API_KEY')
        if not deepgram_key:
            return jsonify({'error': 'DEEPGRAM_API_KEY not configured'}), 500

        resp = http_requests.post(
            f'https://api.deepgram.com/v1/speak?model={voice_id}',
            headers={
                'Authorization': f'Token {deepgram_key}',
                'Content-Type': 'application/json',
            },
            json={"text": text},
            timeout=30,
        )
        resp.raise_for_status()

        return FlaskResponse(resp.content, mimetype='audio/mpeg', headers={
            'Access-Control-Allow-Origin': '*',
            'Cache-Control': 'public, max-age=3600',
        })

    except Exception as e:
        print(f"[preview-voice] Error: {e}")
        return jsonify({'error': str(e)}), 500


def _voiceover_job_worker(job_id, scripts, voice_id, user_id, voice_provider='qwen'):
    """Background worker: generate voiceover audio for each slide, ZIP MP3s, upload."""
    try:
        import requests as http_requests
        from pydub import AudioSegment
        import zipfile
        import io
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import time
        import shutil

        supabase = get_supabase_client()
        alibaba_key = os.environ.get('ALIBABA_API_KEY')
        deepgram_key = os.environ.get('DEEPGRAM_API_KEY')

        # Update job → processing
        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        job_dir = os.path.join(OUTPUT_DIR, f'voiceover_{job_id}')
        os.makedirs(job_dir, exist_ok=True)

        use_deepgram = voice_provider == 'deepgram' and deepgram_key

        def synthesize_chunk_deepgram(text_chunk, attempt=0):
            """Call Deepgram Aura TTS API. Returns MP3 bytes directly."""
            max_attempts = 4
            try:
                resp = http_requests.post(
                    f'https://api.deepgram.com/v1/speak?model={voice_id}',
                    headers={
                        'Authorization': f'Token {deepgram_key}',
                        'Content-Type': 'application/json',
                    },
                    json={"text": text_chunk},
                    timeout=120,
                )
                resp.raise_for_status()
                return resp.content  # MP3 bytes
            except Exception as e:
                if attempt < max_attempts - 1:
                    wait = (2 ** attempt) + 1
                    print(f"[TTS-DG] Chunk retry {attempt+1}/{max_attempts}: {e}, waiting {wait}s")
                    time.sleep(wait)
                    return synthesize_chunk_deepgram(text_chunk, attempt + 1)
                raise

        def synthesize_chunk_qwen(text_chunk, attempt=0):
            """Call Qwen TTS synthesis API for a single chunk. Returns WAV bytes."""
            max_attempts = 4
            try:
                resp = http_requests.post(
                    'https://dashscope-intl.aliyuncs.com/api/v1/services/aigc/multimodal-generation/generation',
                    headers={
                        'Authorization': f'Bearer {alibaba_key}',
                        'Content-Type': 'application/json',
                    },
                    json={
                        "model": "qwen3-tts-vc-2026-01-22",
                        "input": {
                            "text": text_chunk,
                            "voice": voice_id,
                        }
                    },
                    timeout=120,
                )
                resp.raise_for_status()
                result = resp.json()
                audio_url = result.get('output', {}).get('audio', {}).get('url')
                if not audio_url:
                    raise Exception(f'No audio URL in response: {result}')

                # Download the WAV file
                wav_resp = http_requests.get(audio_url, timeout=60)
                wav_resp.raise_for_status()
                return wav_resp.content

            except Exception as e:
                if attempt < max_attempts - 1:
                    wait = (2 ** attempt) + 1
                    print(f"[TTS] Chunk retry {attempt+1}/{max_attempts}: {e}, waiting {wait}s")
                    time.sleep(wait)
                    return synthesize_chunk_qwen(text_chunk, attempt + 1)
                raise

        synthesize_chunk = synthesize_chunk_deepgram if use_deepgram else synthesize_chunk_qwen
        print(f"[TTS] Using {'Deepgram' if use_deepgram else 'Qwen'} provider, voice: {voice_id}")

        slide_mp3_paths = []

        for slide_idx, script in enumerate(scripts):
            script = (script or '').strip()
            if not script:
                # Generate a brief silence for empty scripts
                silence = AudioSegment.silent(duration=1000)
                mp3_path = os.path.join(job_dir, f'slide_{slide_idx + 1:02d}.mp3')
                silence.export(mp3_path, format='mp3')
                slide_mp3_paths.append(mp3_path)
                continue

            # Chunk the script
            chunks = chunk_text_for_tts(script, max_chars=512)
            print(f"[TTS] Slide {slide_idx + 1}: {len(chunks)} chunks, {len(script)} chars")

            # Synthesize chunks (parallelize within each slide)
            chunk_wavs = [None] * len(chunks)

            with ThreadPoolExecutor(max_workers=3) as executor:
                future_to_idx = {
                    executor.submit(synthesize_chunk, chunk): i
                    for i, chunk in enumerate(chunks)
                }
                for future in as_completed(future_to_idx):
                    idx = future_to_idx[future]
                    chunk_wavs[idx] = future.result()

            # Concatenate audio chunks using pydub
            combined = AudioSegment.empty()
            for audio_bytes in chunk_wavs:
                if use_deepgram:
                    segment = AudioSegment.from_mp3(io.BytesIO(audio_bytes))
                else:
                    segment = AudioSegment.from_wav(io.BytesIO(audio_bytes))
                combined += segment

            # Export as MP3
            mp3_path = os.path.join(job_dir, f'slide_{slide_idx + 1:02d}.mp3')
            combined.export(mp3_path, format='mp3', bitrate='128k')
            slide_mp3_paths.append(mp3_path)
            print(f"[TTS] Slide {slide_idx + 1}: MP3 exported ({len(combined)}ms)")

        # ── Create ZIP of all slide MP3s ──
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        zip_filename = f'voiceover_{job_id[:8]}_{timestamp}.zip'
        output_zip_path = os.path.join(OUTPUT_DIR, zip_filename)

        with zipfile.ZipFile(output_zip_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for mp3_path in slide_mp3_paths:
                zf.write(mp3_path, os.path.basename(mp3_path))

        print(f"[TTS] ZIP created: {zip_filename} ({len(slide_mp3_paths)} files)")

        # ── Upload ZIP to Supabase Storage ──
        storage_path = f'voiceovers/{zip_filename}'
        with open(output_zip_path, 'rb') as f:
            file_bytes = f.read()

        supabase.storage.from_('course-files').upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": "application/zip"}
        )
        public_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Save to generated_files table
        file_record = supabase.table('generated_files').insert({
            'file_type': 'voiceover_audio',
            'file_url': public_url,
            'file_size': len(file_bytes),
            'metadata': json.dumps({
                'slide_count': len(scripts),
                'voice_id': voice_id,
            })
        }).execute()

        # Update job → completed
        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Cleanup
        shutil.rmtree(job_dir, ignore_errors=True)
        if os.path.exists(output_zip_path):
            os.remove(output_zip_path)

        print(f"[TTS] Voiceover audio ZIP generated: {zip_filename} ({len(scripts)} slides)")

    except Exception as err:
        import traceback
        print(f"[TTS] Job {job_id} failed: {err}")
        print(f"[TTS] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
        except Exception:
            pass


def _voiceover_video_job_worker(job_id, audio_zip_url, slide_image_urls, user_id):
    """Background worker: download audio ZIP + slide images, render MP4 video."""
    try:
        import requests as http_requests
        from pydub import AudioSegment
        import subprocess
        import zipfile
        import io
        import shutil

        supabase = get_supabase_client()

        # Update job → processing
        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        job_dir = os.path.join(OUTPUT_DIR, f'video_{job_id}')
        os.makedirs(job_dir, exist_ok=True)

        # ── Phase 1: Download and extract audio ZIP ──
        print(f"[MP4] Downloading audio ZIP...")
        zip_resp = http_requests.get(audio_zip_url, timeout=120)
        zip_resp.raise_for_status()

        audio_dir = os.path.join(job_dir, 'audio')
        os.makedirs(audio_dir, exist_ok=True)

        with zipfile.ZipFile(io.BytesIO(zip_resp.content)) as zf:
            zf.extractall(audio_dir)

        # Collect MP3 paths sorted by filename
        mp3_files = sorted([
            os.path.join(audio_dir, f) for f in os.listdir(audio_dir)
            if f.endswith('.mp3')
        ])
        print(f"[MP4] Extracted {len(mp3_files)} MP3 files from ZIP")

        # ── Phase 2: Download slide images ──
        slide_image_paths = []
        for i, url in enumerate(slide_image_urls):
            img_path = os.path.join(job_dir, f'slide_{i + 1:02d}.png')
            img_resp = http_requests.get(url, timeout=30)
            img_resp.raise_for_status()
            with open(img_path, 'wb') as f:
                f.write(img_resp.content)
            slide_image_paths.append(img_path)
            print(f"[MP4] Downloaded slide image {i + 1}/{len(slide_image_urls)}")

        # ── Phase 3: Create per-slide video segments ──
        segment_paths = []
        total_duration_ms = 0

        for i, mp3_path in enumerate(mp3_files):
            slide_audio = AudioSegment.from_mp3(mp3_path)

            # Create padded audio: 1s silence + audio + 2s silence
            silence_1s = AudioSegment.silent(duration=1000)
            silence_2s = AudioSegment.silent(duration=2000)
            padded = silence_1s + slide_audio + silence_2s

            padded_path = os.path.join(job_dir, f'padded_{i + 1:02d}.mp3')
            padded.export(padded_path, format='mp3', bitrate='192k')

            padded_duration_s = len(padded) / 1000.0
            total_duration_ms += len(padded)

            # Use corresponding slide image, or last available image
            if i < len(slide_image_paths):
                img_path = slide_image_paths[i]
            elif slide_image_paths:
                img_path = slide_image_paths[-1]
            else:
                img_path = os.path.join(job_dir, 'blank.png')
                if not os.path.exists(img_path):
                    subprocess.run([
                        'ffmpeg', '-y', '-f', 'lavfi', '-i',
                        'color=c=black:s=1920x1080:d=1',
                        '-frames:v', '1', img_path
                    ], check=True, capture_output=True)

            # Create video segment via ffmpeg
            # -preset ultrafast: dramatically reduces RAM usage (avoids OOM on Railway)
            # -crf 28: slightly higher compression, less memory for rate control
            # -threads 2: cap threads to avoid memory pressure from parallel encoding
            # 1280x720: lower resolution reduces memory by ~55% vs 1920x1080
            segment_path = os.path.join(job_dir, f'segment_{i + 1:02d}.mp4')
            ffmpeg_result = subprocess.run([
                'ffmpeg', '-y',
                '-loop', '1', '-framerate', '1', '-i', img_path,
                '-i', padded_path,
                '-c:v', 'libx264', '-preset', 'ultrafast', '-tune', 'stillimage',
                '-crf', '28', '-threads', '2',
                '-c:a', 'aac', '-b:a', '128k',
                '-t', str(padded_duration_s),
                '-vf', 'scale=1280:720:force_original_aspect_ratio=decrease,pad=1280:720:(ow-iw)/2:(oh-ih)/2:black,format=yuv420p',
                '-movflags', '+faststart',
                '-shortest',
                segment_path
            ], capture_output=True, text=True)
            if ffmpeg_result.returncode != 0:
                err_lines = [l for l in ffmpeg_result.stderr.split('\n') if l.strip() and not l.strip().startswith('frame=')]
                print(f"[MP4] ffmpeg segment error (slide {i+1}, rc={ffmpeg_result.returncode}):\n" + '\n'.join(err_lines[-30:]))
                raise Exception(f"Video encoding failed for slide {i+1} (rc={ffmpeg_result.returncode}). Check server logs for details.")

            segment_paths.append(segment_path)
            print(f"[MP4] Segment {i + 1}/{len(mp3_files)}: {padded_duration_s:.1f}s")

        # ── Phase 4: Concatenate all segments into final MP4 ──
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        mp4_filename = f'voiceover_{job_id[:8]}_{timestamp}.mp4'
        output_mp4_path = os.path.join(OUTPUT_DIR, mp4_filename)

        concat_list_path = os.path.join(job_dir, 'concat.txt')
        with open(concat_list_path, 'w') as f:
            for seg_path in segment_paths:
                f.write(f"file '{seg_path}'\n")

        concat_result = subprocess.run([
            'ffmpeg', '-y',
            '-f', 'concat', '-safe', '0',
            '-i', concat_list_path,
            '-c', 'copy',
            output_mp4_path
        ], capture_output=True, text=True)
        if concat_result.returncode != 0:
            print(f"[MP4] ffmpeg concat error: {concat_result.stderr[-1000:]}")
            raise Exception(f"ffmpeg concat failed: {concat_result.stderr[-300:]}")

        print(f"[MP4] Final video created: {mp4_filename} ({total_duration_ms / 1000:.1f}s total)")

        # ── Phase 5: Upload MP4 to Supabase Storage ──
        storage_path = f'voiceovers/{mp4_filename}'
        with open(output_mp4_path, 'rb') as f:
            file_bytes = f.read()

        supabase.storage.from_('course-files').upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": "video/mp4"}
        )
        public_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Save to generated_files table
        file_record = supabase.table('generated_files').insert({
            'file_type': 'voiceover_video',
            'file_url': public_url,
            'file_size': len(file_bytes),
            'metadata': json.dumps({
                'slide_count': len(mp3_files),
                'duration_ms': total_duration_ms,
                'resolution': '1920x1080',
            })
        }).execute()

        # Update job → completed
        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Cleanup
        shutil.rmtree(job_dir, ignore_errors=True)
        if os.path.exists(output_mp4_path):
            os.remove(output_mp4_path)

        print(f"[MP4] Voiceover video generated: {mp4_filename} ({len(mp3_files)} slides)")

    except Exception as err:
        import traceback
        print(f"[MP4] Job {job_id} failed: {err}")
        print(f"[MP4] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
        except Exception:
            pass


@app.route('/generate-voiceover-audio', methods=['POST'])
def generate_voiceover_audio():
    """
    Start async voiceover audio generation (audio ZIP only).
    Expected JSON: { "scripts": ["slide 1 text", ...], "voice_id": "...", "user_id": "uuid" }
    Returns: { "job_id": "uuid" } immediately; poll /job-status/<job_id>.
    """
    try:
        data = request.json or {}
        scripts = data.get('scripts', [])
        voice_id = data.get('voice_id')
        voice_provider = data.get('voice_provider', 'qwen')
        user_id = data.get('user_id')

        if not scripts:
            return jsonify({'error': 'scripts array is required'}), 400
        if not voice_id:
            return jsonify({'error': 'voice_id is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        supabase = get_supabase_client()

        # Create job row
        job = supabase.table('generation_jobs').insert({
            'job_type': 'voiceover_audio',
            'status': 'pending',
            'user_id': user_id,
        }).execute()

        job_id = job.data[0]['id']

        # Spawn background worker
        threading.Thread(
            target=_voiceover_job_worker,
            args=(job_id, scripts, voice_id, user_id, voice_provider),
            daemon=True
        ).start()

        return jsonify({'success': True, 'job_id': job_id})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/generate-voiceover-video', methods=['POST'])
def generate_voiceover_video():
    """
    Start async voiceover video generation from existing audio ZIP + slide images.
    Expected JSON: { "audio_zip_url": "...", "slide_image_urls": [...], "user_id": "uuid" }
    Returns: { "job_id": "uuid" } immediately; poll /job-status/<job_id>.
    """
    try:
        data = request.json or {}
        audio_zip_url = data.get('audio_zip_url')
        slide_image_urls = data.get('slide_image_urls', [])
        user_id = data.get('user_id')

        if not audio_zip_url:
            return jsonify({'error': 'audio_zip_url is required'}), 400
        if not slide_image_urls:
            return jsonify({'error': 'slide_image_urls array is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        supabase = get_supabase_client()

        # Create job row
        job = supabase.table('generation_jobs').insert({
            'job_type': 'voiceover_video',
            'status': 'pending',
            'user_id': user_id,
        }).execute()

        job_id = job.data[0]['id']

        # Spawn background worker
        threading.Thread(
            target=_voiceover_video_job_worker,
            args=(job_id, audio_zip_url, slide_image_urls, user_id),
            daemon=True
        ).start()

        return jsonify({'success': True, 'job_id': job_id})

    except Exception as e:
        return jsonify({'error': str(e)}), 500


# --- Direct file-return endpoints (kept for local/direct use) ---

@app.route('/generate-powerpoint', methods=['POST'])
def generate_ppt():
    """Generate PowerPoint directly and return file (no Supabase)"""
    try:
        data = request.json
        if not data or 'slides' not in data:
            return jsonify({'error': 'Missing course outline data'}), 400

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"course_{timestamp}.pptx"
        filepath = os.path.join(OUTPUT_DIR, filename)
        generate_powerpoint_file(data, filepath)

        return send_file(
            filepath,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/generate-scorm', methods=['POST'])
def generate_scorm():
    """Generate SCORM package directly and return file (no Supabase)"""
    try:
        data = request.json
        if not data or 'slides' not in data:
            return jsonify({'error': 'Missing course outline data'}), 400

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"course_{timestamp}.zip"
        filepath = os.path.join(OUTPUT_DIR, filename)
        generate_scorm_package(data, filepath)

        return send_file(
            filepath,
            mimetype='application/zip',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/generate-scorm-from-video', methods=['POST'])
def generate_scorm_from_video():
    """Generate SCORM 1.2 package wrapping an MP4 video."""
    try:
        data = request.json
        title = data.get('title', 'Training Video')
        video_url = data.get('video_url')
        user_id = data.get('user_id')

        if not video_url:
            return jsonify({'error': 'video_url is required'}), 400
        if not user_id:
            return jsonify({'error': 'user_id is required'}), 400

        from generate_scorm import generate_video_scorm_package

        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"scorm_video_{timestamp}.zip"
        filepath = os.path.join(OUTPUT_DIR, filename)

        generate_video_scorm_package(title, video_url, filepath)

        # Upload to Supabase Storage
        supabase = get_supabase_client()
        storage_path = f"scorm/{user_id}_{timestamp}.zip"
        with open(filepath, 'rb') as f:
            supabase.storage.from_('course-files').upload(
                path=storage_path, file=f.read(),
                file_options={"content-type": "application/zip"}
            )
        download_url = supabase.storage.from_('course-files').get_public_url(storage_path)

        # Clean up local file
        try:
            os.remove(filepath)
        except Exception:
            pass

        return jsonify({'success': True, 'download_url': download_url, 'filename': filename})

    except Exception as e:
        print(f"[SCORM-Video] Error: {e}")
        return jsonify({'error': str(e)}), 500


@app.route('/search-chunks', methods=['POST'])
def search_chunks():
    """
    Embed a question and return the top-k most semantically similar document chunks.

    Expected JSON:
      {
        "question":     str,
        "client_id":    str (UUID),
        "document_ids": [str, ...],   # optional filter; empty/omitted = all docs
        "top_k":        int           # default 8
      }
    """
    try:
        data         = request.json or {}
        question     = data.get('question', '').strip()
        document_ids = data.get('document_ids') or []
        client_id    = data.get('client_id', '').strip()
        top_k        = int(data.get('top_k', 8))

        if not question:
            return jsonify({'error': 'question is required'}), 400
        if not client_id:
            return jsonify({'error': 'client_id is required'}), 400

        from embeddings import embed_texts
        query_embedding = embed_texts([question])[0]

        supabase = get_supabase_client()
        result = supabase.rpc('match_chunks', {
            'query_embedding': query_embedding,
            'match_document_ids': document_ids if document_ids else None,
            'match_client_id': client_id,
            'match_count': top_k,
        }).execute()

        return jsonify({'success': True, 'chunks': result.data or []})

    except Exception as e:
        print(f"[search-chunks] Error: {e}")
        return jsonify({'error': str(e)}), 500


def _reindex_worker(client_id):
    """Background thread: chunk + embed all documents for a user."""
    import gc
    try:
        from embeddings import embed_texts, chunk_text
        supabase = get_supabase_client()

        result = supabase.table('documents') \
            .select('id, filename, extracted_text') \
            .eq('client_id', client_id) \
            .execute()

        documents = result.data or []
        print(f"[reindex] Starting reindex for {client_id}: {len(documents)} documents")

        for doc in documents:
            doc_id         = doc['id']
            filename       = doc['filename']
            extracted_text = doc.get('extracted_text') or ''

            if not extracted_text.strip():
                print(f"[reindex] Skipping {filename} — no extracted text")
                continue

            try:
                supabase.table('document_chunks').delete().eq('document_id', doc_id).execute()

                chunks = chunk_text(extracted_text)

                # Embed in small batches of 20 to avoid OOM spikes on Railway
                EMBED_BATCH = 20
                all_embeddings = []
                for b in range(0, len(chunks), EMBED_BATCH):
                    batch_embs = embed_texts(chunks[b:b + EMBED_BATCH])
                    all_embeddings.extend(batch_embs)
                    gc.collect()  # free intermediate tensors between batches

                chunk_rows = [
                    {
                        'document_id': doc_id,
                        'client_id':   client_id,
                        'chunk_index': i,
                        'chunk_text':  c,
                        'embedding':   e,
                    }
                    for i, (c, e) in enumerate(zip(chunks, all_embeddings))
                ]

                # Insert in batches of 50 (smaller payload to Supabase)
                for batch_start in range(0, len(chunk_rows), 50):
                    supabase.table('document_chunks').insert(
                        chunk_rows[batch_start:batch_start + 50]
                    ).execute()

                print(f"[reindex] ✓ {filename}: {len(chunks)} chunks")

            except Exception as doc_err:
                print(f"[reindex] ✗ {filename}: {doc_err}")

            finally:
                gc.collect()  # free memory between documents

        print(f"[reindex] Done for {client_id}")

    except Exception as e:
        print(f"[reindex] Fatal error: {e}")


@app.route('/reindex-documents', methods=['POST'])
def reindex_documents():
    """
    Re-chunk and re-embed all existing documents for a user.
    Runs in a background thread and returns immediately — check Railway logs for progress.

    Expected JSON: { "client_id": str (UUID) }
    """
    try:
        data      = request.json or {}
        client_id = data.get('client_id', '').strip()

        if not client_id:
            return jsonify({'error': 'client_id is required'}), 400

        # Kick off background thread and return immediately (avoids Railway timeout)
        import threading
        t = threading.Thread(target=_reindex_worker, args=(client_id,), daemon=True)
        t.start()

        return jsonify({
            'success': True,
            'message': 'Reindexing started in background. Check Railway logs for progress.',
            'client_id': client_id,
        })

    except Exception as e:
        print(f"[reindex-documents] Error: {e}")
        return jsonify({'error': str(e)}), 500


###############################################################################
# SCORM Cloud LMS Push
###############################################################################

@app.route('/extract-scorm-cloud-text', methods=['POST'])
def extract_scorm_cloud_text():
    """Download a SCORM package from SCORM Cloud and extract text content."""
    try:
        import requests as http_requests
        import zipfile
        import io
        import xml.etree.ElementTree as ET
        from bs4 import BeautifulSoup

        data = request.get_json()
        course_id = data.get('course_id')
        app_id = data.get('app_id')
        secret_key = data.get('secret_key')

        if not course_id:
            return jsonify({'error': 'course_id is required'}), 400
        if not app_id or not secret_key:
            return jsonify({'error': 'SCORM Cloud credentials are required'}), 400

        # 1. Download the course ZIP from SCORM Cloud
        print(f"[extract-scorm] Downloading ZIP for courseId={course_id}")
        zip_resp = http_requests.get(
            f"https://cloud.scorm.com/api/v2/courses/{course_id}/zip",
            auth=(app_id, secret_key),
            timeout=120,
        )
        if zip_resp.status_code == 404:
            return jsonify({'error': 'Course not found in SCORM Cloud'}), 404
        if zip_resp.status_code == 401:
            return jsonify({'error': 'SCORM Cloud authentication failed'}), 401
        if zip_resp.status_code != 200:
            return jsonify({'error': f'Failed to download course (HTTP {zip_resp.status_code})'}), 500

        zip_bytes = zip_resp.content
        print(f"[extract-scorm] Downloaded {len(zip_bytes)} bytes")

        # 2. Open ZIP and find content files
        zf = zipfile.ZipFile(io.BytesIO(zip_bytes))
        filenames = zf.namelist()

        # Try to find imsmanifest.xml to identify the correct HTML entry files
        html_files = []
        course_title = course_id  # fallback title

        manifest_names = [f for f in filenames if f.lower().endswith('imsmanifest.xml')]
        if manifest_names:
            try:
                manifest_content = zf.read(manifest_names[0]).decode('utf-8', errors='ignore')
                root = ET.fromstring(manifest_content)

                # Extract course title from manifest metadata
                # Try multiple namespace patterns
                for ns_prefix in [
                    '{http://www.imsproject.org/xsd/imscp_rootv1p1p2}',
                    '{http://www.imsglobal.org/xsd/imscp_v1p1}',
                    '',
                ]:
                    title_el = root.find(f'.//{ns_prefix}title')
                    if title_el is not None and title_el.text:
                        course_title = title_el.text.strip()
                        break

                # Find resource hrefs (HTML entry points)
                for ns_prefix in [
                    '{http://www.imsproject.org/xsd/imscp_rootv1p1p2}',
                    '{http://www.imsglobal.org/xsd/imscp_v1p1}',
                    '',
                ]:
                    resources = root.findall(f'.//{ns_prefix}resource')
                    if resources:
                        for res in resources:
                            href = res.get('href', '')
                            if href and (href.endswith('.html') or href.endswith('.htm')):
                                html_files.append(href)
                        break

                print(f"[extract-scorm] Manifest title: {course_title}, {len(html_files)} HTML resources")
            except Exception as e:
                print(f"[extract-scorm] Warning: failed to parse manifest: {e}")

        # If manifest parsing didn't find HTML files, find them in the ZIP directly
        if not html_files:
            html_files = [f for f in filenames
                         if (f.endswith('.html') or f.endswith('.htm'))
                         and not f.startswith('__MACOSX')
                         and '/.' not in f]

        if not html_files:
            return jsonify({'error': 'No HTML content found in SCORM package'}), 400

        print(f"[extract-scorm] Found {len(html_files)} HTML files to parse")

        # 3. Extract text from HTML files
        slides = []
        slide_num = 0

        for html_file in html_files:
            try:
                html_content = zf.read(html_file).decode('utf-8', errors='ignore')
                soup = BeautifulSoup(html_content, 'html.parser')

                # Remove script, style, noscript tags
                for tag in soup(['script', 'style', 'noscript', 'nav', 'header', 'footer']):
                    tag.decompose()

                # Strategy 1: Look for slide-like div structures (our own SCORM packages)
                slide_divs = soup.find_all('div', class_=lambda c: c and 'slide' in str(c).lower()) if soup.find('div', class_=lambda c: c and 'slide' in str(c).lower()) else []

                if slide_divs:
                    for div in slide_divs:
                        slide_num += 1
                        # Find title (h1, h2, h3, or first strong/b)
                        title_el = div.find(['h1', 'h2', 'h3'])
                        title = title_el.get_text(strip=True) if title_el else f'Slide {slide_num}'
                        if title_el:
                            title_el.decompose()
                        text = div.get_text(separator='\n', strip=True)
                        if text.strip():
                            slides.append({
                                'number': slide_num,
                                'title': title[:200],
                                'text': text[:3000]
                            })
                else:
                    # Strategy 2: Treat each HTML file as a page, extract sections
                    # Find the main content area (body or main content div)
                    body = soup.find('body') or soup

                    # Try to split by headings
                    headings = body.find_all(['h1', 'h2', 'h3'])
                    if headings:
                        for heading in headings:
                            slide_num += 1
                            title = heading.get_text(strip=True)

                            # Collect text between this heading and the next
                            text_parts = []
                            sibling = heading.find_next_sibling()
                            while sibling and sibling.name not in ['h1', 'h2', 'h3']:
                                t = sibling.get_text(separator='\n', strip=True)
                                if t:
                                    text_parts.append(t)
                                sibling = sibling.find_next_sibling()

                            text = '\n'.join(text_parts)
                            if title.strip() or text.strip():
                                slides.append({
                                    'number': slide_num,
                                    'title': title[:200] or f'Section {slide_num}',
                                    'text': text[:3000]
                                })
                    else:
                        # No headings — treat entire page as one slide
                        text = body.get_text(separator='\n', strip=True)
                        if text.strip() and len(text.strip()) > 20:
                            slide_num += 1
                            # Use filename as title
                            page_title = html_file.rsplit('/', 1)[-1].replace('.html', '').replace('.htm', '').replace('_', ' ').replace('-', ' ').title()
                            slides.append({
                                'number': slide_num,
                                'title': page_title[:200],
                                'text': text[:3000]
                            })
            except Exception as e:
                print(f"[extract-scorm] Warning: failed to parse {html_file}: {e}")
                continue

        zf.close()

        if not slides:
            return jsonify({'error': 'Could not extract text content from SCORM package. The content may be JavaScript-rendered or image-based.'}), 400

        print(f"[extract-scorm] Extracted {len(slides)} slides/sections from {course_id}")

        return jsonify({
            'success': True,
            'slides': slides,
            'course_title': course_title,
            'slide_count': len(slides),
        })

    except Exception as e:
        print(f"[extract-scorm] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


def _transcription_job_worker(job_id, media_files, course_id, user_id):
    """Background worker: transcribe audio/video files from a SCORM package using DashScope STT."""
    try:
        import requests as http_requests
        import base64
        import re
        import time

        supabase = get_supabase_client()
        alibaba_key = os.environ.get('ALIBABA_API_KEY')

        # Update job → processing
        supabase.table('generation_jobs').update({
            'status': 'processing',
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        transcriptions = []

        for media_file in media_files:
            filename = media_file['filename']
            audio_bytes = media_file['data']
            print(f"[transcribe] Processing {filename} ({len(audio_bytes)} bytes)")

            try:
                # For video files, extract audio track using ffmpeg
                ext = filename.rsplit('.', 1)[-1].lower()
                if ext in ('mp4', 'webm', 'mov', 'avi'):
                    import subprocess
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix=f'.{ext}', delete=False) as vf:
                        vf.write(audio_bytes)
                        video_path = vf.name
                    audio_path = video_path + '.wav'
                    try:
                        subprocess.run([
                            'ffmpeg', '-i', video_path, '-vn', '-acodec', 'pcm_s16le',
                            '-ar', '16000', '-ac', '1', audio_path, '-y'
                        ], capture_output=True, timeout=120, check=True)
                        with open(audio_path, 'rb') as af:
                            audio_bytes = af.read()
                    finally:
                        for p in [video_path, audio_path]:
                            if os.path.exists(p):
                                os.remove(p)

                # Base64-encode audio for DashScope API
                audio_b64 = base64.b64encode(audio_bytes).decode('utf-8')

                # Determine audio format for the API
                if ext in ('mp4', 'webm', 'mov', 'avi'):
                    audio_format = 'wav'  # we converted to wav above
                elif ext == 'mp3':
                    audio_format = 'mp3'
                elif ext in ('wav',):
                    audio_format = 'wav'
                elif ext == 'ogg':
                    audio_format = 'ogg'
                elif ext == 'm4a':
                    audio_format = 'mp4'
                else:
                    audio_format = 'wav'

                # Call DashScope Paraformer STT API
                resp = http_requests.post(
                    'https://dashscope-intl.aliyuncs.com/api/v1/services/audio/asr/transcription',
                    headers={
                        'Authorization': f'Bearer {alibaba_key}',
                        'Content-Type': 'application/json',
                    },
                    json={
                        "model": "paraformer-v2",
                        "input": {
                            "audio": f"data:audio/{audio_format};base64,{audio_b64}",
                        },
                        "parameters": {
                            "language": "en",
                        }
                    },
                    timeout=180,
                )
                resp.raise_for_status()
                result = resp.json()

                # Extract transcript text from response
                transcript = ''
                output = result.get('output', {})
                if isinstance(output, dict):
                    # Try different response formats
                    if 'sentence' in output:
                        sentences = output['sentence']
                        if isinstance(sentences, list):
                            transcript = ' '.join(s.get('text', '') for s in sentences)
                        elif isinstance(sentences, dict):
                            transcript = sentences.get('text', '')
                    elif 'text' in output:
                        transcript = output['text']
                    elif 'results' in output:
                        for r in output['results']:
                            transcript += r.get('text', '') + ' '

                transcript = transcript.strip()
                if not transcript:
                    print(f"[transcribe] Warning: no transcript returned for {filename}")
                    continue

                # Try to associate with a slide number from filename
                associated_slide = None
                slide_match = re.search(r'(?:slide|s)[\s_-]*(\d+)', filename, re.IGNORECASE)
                if slide_match:
                    associated_slide = int(slide_match.group(1))

                # Estimate duration from audio (rough: 16000 samples/sec for WAV)
                duration_seconds = len(audio_bytes) / 32000  # rough estimate

                transcriptions.append({
                    'filename': filename,
                    'transcript': transcript[:10000],  # cap at 10k chars per file
                    'duration_seconds': round(duration_seconds, 1),
                    'associated_slide': associated_slide,
                })

                print(f"[transcribe] {filename}: {len(transcript)} chars transcribed")

            except Exception as file_err:
                print(f"[transcribe] Warning: failed to transcribe {filename}: {file_err}")
                continue

        # Store transcriptions in generated_files metadata
        file_record = supabase.table('generated_files').insert({
            'file_type': 'scorm_transcription',
            'file_url': '',  # no file URL, data is in metadata
            'file_size': 0,
            'metadata': json.dumps({
                'transcriptions': transcriptions,
                'course_id': course_id,
                'media_files_processed': len(media_files),
            })
        }).execute()

        # Update job → completed
        supabase.table('generation_jobs').update({
            'status': 'completed',
            'result_file_id': file_record.data[0]['id'],
            'updated_at': datetime.now().isoformat()
        }).eq('id', job_id).execute()

        # Also cache in lms_course_transcriptions for the edge function
        try:
            supabase.table('lms_course_transcriptions').upsert({
                'user_id': user_id,
                'course_id': course_id,
                'transcriptions': transcriptions,
            }, on_conflict='user_id,course_id').execute()
        except Exception as cache_err:
            print(f"[transcribe] Warning: failed to cache transcriptions: {cache_err}")

        print(f"[transcribe] Job {job_id} completed: {len(transcriptions)} files transcribed")

    except Exception as err:
        import traceback
        print(f"[transcribe] Job {job_id} failed: {err}")
        print(f"[transcribe] Traceback:\n{traceback.format_exc()}")
        try:
            supabase = get_supabase_client()
            supabase.table('generation_jobs').update({
                'status': 'failed',
                'error_message': str(err)[:500],
                'updated_at': datetime.now().isoformat()
            }).eq('id', job_id).execute()
        except Exception:
            pass


@app.route('/extract-scorm-cloud-media', methods=['POST'])
def extract_scorm_cloud_media():
    """Download a SCORM package from SCORM Cloud and transcribe audio/video content."""
    try:
        import requests as http_requests
        import zipfile
        import io

        data = request.get_json()
        course_id = data.get('course_id')
        app_id = data.get('app_id')
        secret_key = data.get('secret_key')
        user_id = data.get('user_id')

        if not course_id:
            return jsonify({'error': 'course_id is required'}), 400
        if not app_id or not secret_key:
            return jsonify({'error': 'SCORM Cloud credentials are required'}), 400

        alibaba_key = os.environ.get('ALIBABA_API_KEY')
        if not alibaba_key:
            return jsonify({'error': 'ALIBABA_API_KEY not configured (needed for transcription)'}), 500

        # 1. Download the course ZIP from SCORM Cloud
        print(f"[extract-scorm-media] Downloading ZIP for courseId={course_id}")
        zip_resp = http_requests.get(
            f"https://cloud.scorm.com/api/v2/courses/{course_id}/zip",
            auth=(app_id, secret_key),
            timeout=120,
        )
        if zip_resp.status_code == 404:
            return jsonify({'error': 'Course not found in SCORM Cloud'}), 404
        if zip_resp.status_code == 401:
            return jsonify({'error': 'SCORM Cloud authentication failed'}), 401
        if zip_resp.status_code != 200:
            return jsonify({'error': f'Failed to download course (HTTP {zip_resp.status_code})'}), 500

        zip_bytes = zip_resp.content
        print(f"[extract-scorm-media] Downloaded {len(zip_bytes)} bytes")

        # 2. Scan ZIP for media files
        MEDIA_EXTENSIONS = {'.mp3', '.wav', '.ogg', '.m4a', '.mp4', '.webm', '.mov', '.avi'}
        MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB per file

        zf = zipfile.ZipFile(io.BytesIO(zip_bytes))
        media_files = []

        for info in zf.infolist():
            ext = ('.' + info.filename.rsplit('.', 1)[-1].lower()) if '.' in info.filename else ''
            if ext in MEDIA_EXTENSIONS and info.file_size < MAX_FILE_SIZE and not info.filename.startswith('__MACOSX'):
                media_files.append({
                    'filename': info.filename.rsplit('/', 1)[-1],  # just the filename, not path
                    'data': zf.read(info.filename),
                })

        zf.close()

        print(f"[extract-scorm-media] Found {len(media_files)} media files in {course_id}")

        # 3. No media? Return immediately
        if len(media_files) == 0:
            return jsonify({
                'success': True,
                'media_files_found': 0,
                'job_id': None,
            })

        # 4. Create async job and spawn worker
        supabase = get_supabase_client()
        job = supabase.table('generation_jobs').insert({
            'job_type': 'scorm_transcription',
            'status': 'pending',
            'user_id': user_id,
        }).execute()

        job_id = job.data[0]['id']

        threading.Thread(
            target=_transcription_job_worker,
            args=(job_id, media_files, course_id, user_id),
            daemon=True
        ).start()

        return jsonify({
            'success': True,
            'job_id': job_id,
            'media_files_found': len(media_files),
        })

    except Exception as e:
        print(f"[extract-scorm-media] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/push-to-scorm-cloud', methods=['POST'])
def push_to_scorm_cloud():
    """Download a SCORM ZIP from storage and upload it to SCORM Cloud via API v2."""
    try:
        data = request.get_json()
        scorm_url = data.get('scorm_url')
        course_title = data.get('course_title', 'Untitled Course')
        app_id = data.get('app_id')
        secret_key = data.get('secret_key')

        if not scorm_url:
            return jsonify({'error': 'scorm_url is required'}), 400
        if not app_id or not secret_key:
            return jsonify({'error': 'SCORM Cloud credentials (app_id, secret_key) are required'}), 400

        import re, time as _time
        import requests as http_requests

        # 1. Download the SCORM ZIP from Supabase Storage
        print(f"[push-to-scorm-cloud] Downloading SCORM from: {scorm_url[:80]}...")
        dl_resp = http_requests.get(scorm_url, timeout=60)
        if dl_resp.status_code != 200:
            return jsonify({'error': f'Failed to download SCORM package (HTTP {dl_resp.status_code})'}), 500
        content_type = dl_resp.headers.get('Content-Type', '')
        if 'zip' not in content_type and 'octet-stream' not in content_type:
            print(f"[push-to-scorm-cloud] Unexpected Content-Type: {content_type}")
            return jsonify({'error': f'Downloaded file is not a valid ZIP (Content-Type: {content_type})'}), 500
        zip_bytes = dl_resp.content
        if len(zip_bytes) < 100:
            return jsonify({'error': 'Downloaded SCORM package is too small — the file may be corrupted or expired'}), 500
        print(f"[push-to-scorm-cloud] Downloaded {len(zip_bytes)} bytes")

        # 2. Generate a courseId (alphanumeric + hyphens, max 200 chars)
        slug = re.sub(r'[^a-zA-Z0-9]+', '-', course_title).strip('-').lower()[:80]
        course_id = f"{slug}-{int(_time.time())}"

        # 3. Upload to SCORM Cloud API v2
        scorm_cloud_url = f"https://cloud.scorm.com/api/v2/courses/importJobs/upload?courseId={course_id}"
        print(f"[push-to-scorm-cloud] Uploading to SCORM Cloud as courseId={course_id}")

        upload_resp = http_requests.post(
            scorm_cloud_url,
            auth=(app_id, secret_key),
            files={'file': ('course.zip', zip_bytes, 'application/zip')},
            timeout=120,
        )

        if upload_resp.status_code not in (200, 201):
            error_text = upload_resp.text[:500]
            print(f"[push-to-scorm-cloud] Upload failed: {upload_resp.status_code} {error_text}")
            if upload_resp.status_code == 401:
                return jsonify({'error': 'SCORM Cloud authentication failed — check your App ID and Secret Key'}), 401
            return jsonify({'error': f'SCORM Cloud upload failed ({upload_resp.status_code}): {error_text}'}), 500

        upload_result = upload_resp.json()
        import_job_id = upload_result.get('result') or upload_result.get('id') or upload_result.get('jobId')
        print(f"[push-to-scorm-cloud] Import job started: {import_job_id}")

        # 4. Poll for import completion (max 60s, every 3s)
        if not import_job_id:
            print("[push-to-scorm-cloud] No import_job_id returned — skipping poll")
            return jsonify({
                'success': True,
                'course_id': course_id,
                'import_job_id': None,
                'status': 'UNKNOWN',
                'message': 'Upload accepted but no import job ID returned. Check your LMS.',
            })

        status = 'RUNNING'
        poll_failures = 0
        for attempt in range(20):
            _time.sleep(3)
            try:
                poll_resp = http_requests.get(
                    f"https://cloud.scorm.com/api/v2/courses/importJobs/{import_job_id}",
                    auth=(app_id, secret_key),
                    timeout=30,
                )
                if poll_resp.status_code == 200:
                    poll_data = poll_resp.json()
                    status = poll_data.get('status', 'UNKNOWN')
                    print(f"[push-to-scorm-cloud] Poll {attempt+1}: {status}")
                    if status in ('COMPLETE', 'ERROR'):
                        break
                else:
                    poll_failures += 1
                    print(f"[push-to-scorm-cloud] Poll {attempt+1} failed: HTTP {poll_resp.status_code}")
                    if poll_failures >= 3:
                        print("[push-to-scorm-cloud] Too many poll failures, stopping")
                        break
            except Exception as poll_err:
                poll_failures += 1
                print(f"[push-to-scorm-cloud] Poll {attempt+1} exception: {poll_err}")
                if poll_failures >= 3:
                    break

        if status == 'ERROR':
            return jsonify({'error': 'SCORM Cloud import failed — the package may be invalid'}), 500

        result = {
            'success': True,
            'course_id': course_id,
            'import_job_id': import_job_id,
            'status': status,
        }
        if status == 'RUNNING':
            result['message'] = 'Import started but not yet confirmed. Check your LMS.'
        print(f"[push-to-scorm-cloud] Done! courseId={course_id}, status={status}")
        return jsonify(result)

    except Exception as e:
        print(f"[push-to-scorm-cloud] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


def _extract_slides_from_scorm_zip(zip_bytes, fallback_title='Course'):
    """Shared helper: extract text slides from a SCORM ZIP file (bytes). Returns dict with slides, course_title, slide_count."""
    import zipfile
    import io
    import xml.etree.ElementTree as ET
    from bs4 import BeautifulSoup

    zf = zipfile.ZipFile(io.BytesIO(zip_bytes))
    filenames = zf.namelist()
    html_files = []
    course_title = fallback_title

    manifest_names = [f for f in filenames if f.lower().endswith('imsmanifest.xml')]
    if manifest_names:
        try:
            manifest_content = zf.read(manifest_names[0]).decode('utf-8', errors='ignore')
            root = ET.fromstring(manifest_content)
            for ns_prefix in ['{http://www.imsproject.org/xsd/imscp_rootv1p1p2}', '{http://www.imsglobal.org/xsd/imscp_v1p1}', '']:
                title_el = root.find(f'.//{ns_prefix}title')
                if title_el is not None and title_el.text:
                    course_title = title_el.text.strip()
                    break
            for ns_prefix in ['{http://www.imsproject.org/xsd/imscp_rootv1p1p2}', '{http://www.imsglobal.org/xsd/imscp_v1p1}', '']:
                resources = root.findall(f'.//{ns_prefix}resource')
                if resources:
                    for res in resources:
                        href = res.get('href', '')
                        if href and (href.endswith('.html') or href.endswith('.htm')):
                            html_files.append(href)
                    break
        except Exception as e:
            print(f"[extract-scorm] Warning: failed to parse manifest: {e}")

    if not html_files:
        html_files = [f for f in filenames if (f.endswith('.html') or f.endswith('.htm')) and not f.startswith('__MACOSX') and '/.' not in f]

    slides = []
    slide_num = 0
    for html_file in html_files:
        try:
            html_content = zf.read(html_file).decode('utf-8', errors='ignore')
            soup = BeautifulSoup(html_content, 'html.parser')
            for tag in soup(['script', 'style', 'noscript', 'nav', 'header', 'footer']):
                tag.decompose()
            slide_divs = soup.find_all('div', class_=lambda c: c and 'slide' in str(c).lower()) if soup.find('div', class_=lambda c: c and 'slide' in str(c).lower()) else []
            if slide_divs:
                for div in slide_divs:
                    slide_num += 1
                    title_el = div.find(['h1', 'h2', 'h3'])
                    title = title_el.get_text(strip=True) if title_el else f'Slide {slide_num}'
                    if title_el: title_el.decompose()
                    text = div.get_text(separator='\n', strip=True)
                    if text.strip():
                        slides.append({'number': slide_num, 'title': title[:200], 'text': text[:3000]})
            else:
                body = soup.find('body') or soup
                headings = body.find_all(['h1', 'h2', 'h3'])
                if headings:
                    for heading in headings:
                        slide_num += 1
                        title = heading.get_text(strip=True)
                        text_parts = []
                        sibling = heading.find_next_sibling()
                        while sibling and sibling.name not in ['h1', 'h2', 'h3']:
                            t = sibling.get_text(separator='\n', strip=True)
                            if t: text_parts.append(t)
                            sibling = sibling.find_next_sibling()
                        text = '\n'.join(text_parts)
                        if title.strip() or text.strip():
                            slides.append({'number': slide_num, 'title': title[:200] or f'Section {slide_num}', 'text': text[:3000]})
                else:
                    text = body.get_text(separator='\n', strip=True)
                    if text.strip() and len(text.strip()) > 20:
                        slide_num += 1
                        page_title = html_file.rsplit('/', 1)[-1].replace('.html', '').replace('.htm', '').replace('_', ' ').replace('-', ' ').title()
                        slides.append({'number': slide_num, 'title': page_title[:200], 'text': text[:3000]})
        except Exception as e:
            print(f"[extract-scorm] Warning: failed to parse {html_file}: {e}")
    zf.close()
    return {'slides': slides, 'course_title': course_title, 'slide_count': len(slides)}


def _moodle_api_call(moodle_url, moodle_token, wsfunction, params=None):
    """Call a Moodle REST web service function and return the parsed JSON response."""
    import requests as http_requests
    base = moodle_url.rstrip('/')
    url = f"{base}/webservice/rest/server.php"
    payload = {'wstoken': moodle_token, 'moodlewsrestformat': 'json', 'wsfunction': wsfunction}
    if params:
        payload.update(params)
    resp = http_requests.post(url, data=payload, timeout=60)
    resp.raise_for_status()
    data = resp.json()
    if isinstance(data, dict) and data.get('exception'):
        raise Exception(data.get('message', data.get('errorcode', 'Moodle API error')))
    return data


@app.route('/push-to-moodle', methods=['POST'])
def push_to_moodle():
    """Upload a SCORM ZIP to Moodle: upload to draft area, create course, return URL."""
    try:
        import requests as http_requests
        import re

        data = request.get_json()
        scorm_url = data.get('scorm_url')
        course_title = data.get('course_title', 'Untitled Course')
        moodle_url = data.get('moodle_url')
        moodle_token = data.get('moodle_token')

        if not scorm_url:
            return jsonify({'error': 'scorm_url is required'}), 400
        if not moodle_url or not moodle_token:
            return jsonify({'error': 'Moodle credentials are required'}), 400

        base = moodle_url.rstrip('/')

        # 1. Download SCORM ZIP
        print(f"[push-to-moodle] Downloading SCORM ZIP from {scorm_url[:80]}...")
        dl_resp = http_requests.get(scorm_url, timeout=120)
        if dl_resp.status_code != 200:
            return jsonify({'error': f'Failed to download SCORM ZIP (HTTP {dl_resp.status_code})'}), 500
        zip_bytes = dl_resp.content
        print(f"[push-to-moodle] Downloaded {len(zip_bytes)} bytes")

        # 2. Upload ZIP to Moodle draft area
        upload_url = f"{base}/webservice/upload.php?token={moodle_token}"
        slug = re.sub(r'[^a-z0-9]+', '-', course_title.lower()).strip('-')[:50]
        filename = f"{slug}.zip"
        upload_resp = http_requests.post(upload_url, files={'file_1': (filename, zip_bytes, 'application/zip')}, timeout=120)
        if upload_resp.status_code != 200:
            return jsonify({'error': f'Moodle file upload failed (HTTP {upload_resp.status_code})'}), 500
        upload_data = upload_resp.json()
        if isinstance(upload_data, list) and len(upload_data) > 0:
            item_id = upload_data[0].get('itemid')
        elif isinstance(upload_data, dict):
            item_id = upload_data.get('itemid')
        else:
            return jsonify({'error': 'Unexpected Moodle upload response'}), 500

        if not item_id:
            return jsonify({'error': 'Failed to get itemid from Moodle upload'}), 500

        print(f"[push-to-moodle] Uploaded to draft area, itemid={item_id}")

        # 3. Create a new course in Moodle
        shortname = f"{slug}-{int(__import__('time').time())}"
        course_data = _moodle_api_call(moodle_url, moodle_token, 'core_course_create_courses', {
            'courses[0][fullname]': course_title,
            'courses[0][shortname]': shortname,
            'courses[0][categoryid]': 1,
            'courses[0][summary]': f'SCORM course: {course_title}',
            'courses[0][format]': 'singleactivity',
            'courses[0][numsections]': 0,
        })

        if isinstance(course_data, list) and len(course_data) > 0:
            new_course_id = course_data[0].get('id')
        else:
            return jsonify({'error': 'Failed to create Moodle course'}), 500

        print(f"[push-to-moodle] Created course id={new_course_id}")

        course_url = f"{base}/course/view.php?id={new_course_id}"

        return jsonify({
            'success': True,
            'course_id': str(new_course_id),
            'import_job_id': str(item_id),
            'status': 'COMPLETE',
            'course_url': course_url,
            'message': f'Course created. SCORM package uploaded to draft area (itemid={item_id}). Add as SCORM activity in the course to complete setup.',
        })

    except Exception as e:
        print(f"[push-to-moodle] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/extract-moodle-course-content', methods=['POST'])
def extract_moodle_course_content():
    """Fetch course content from Moodle and extract text from SCORM packages."""
    try:
        import requests as http_requests

        data = request.get_json()
        course_id = data.get('course_id')
        moodle_url = data.get('moodle_url')
        moodle_token = data.get('moodle_token')

        if not course_id or not moodle_url or not moodle_token:
            return jsonify({'error': 'course_id, moodle_url, and moodle_token are required'}), 400

        base = moodle_url.rstrip('/')

        # Get course contents (sections + modules)
        contents = _moodle_api_call(moodle_url, moodle_token, 'core_course_get_contents', {'courseid': course_id})

        # Find SCORM modules and their packages
        all_slides = []
        course_title = f'Moodle Course {course_id}'

        # Also try to get the course title
        try:
            courses = _moodle_api_call(moodle_url, moodle_token, 'core_course_get_courses', {'options[ids][0]': course_id})
            if isinstance(courses, list) and courses:
                course_title = courses[0].get('fullname', course_title)
        except Exception:
            pass

        for section in (contents if isinstance(contents, list) else []):
            for module in section.get('modules', []):
                if module.get('modname') == 'scorm':
                    # Try to download SCORM package
                    for content_file in module.get('contents', []):
                        file_url = content_file.get('fileurl')
                        if file_url and file_url.endswith('.zip'):
                            try:
                                dl_url = f"{file_url}&token={moodle_token}" if '?' in file_url else f"{file_url}?token={moodle_token}"
                                resp = http_requests.get(dl_url, timeout=120)
                                if resp.status_code == 200:
                                    result = _extract_slides_from_scorm_zip(resp.content, module.get('name', 'SCORM'))
                                    all_slides.extend(result['slides'])
                            except Exception as e:
                                print(f"[extract-moodle] Warning: failed to download SCORM from {file_url}: {e}")

                # Also extract text from label, page, and other content modules
                elif module.get('modname') in ('label', 'page'):
                    desc = module.get('description', '') or ''
                    if desc.strip():
                        from bs4 import BeautifulSoup
                        soup = BeautifulSoup(desc, 'html.parser')
                        text = soup.get_text(separator='\n', strip=True)
                        if text and len(text) > 20:
                            all_slides.append({
                                'number': len(all_slides) + 1,
                                'title': module.get('name', f'Section {len(all_slides) + 1}')[:200],
                                'text': text[:3000]
                            })

        # Re-number slides
        for i, s in enumerate(all_slides):
            s['number'] = i + 1

        if not all_slides:
            return jsonify({'error': 'No text content found in Moodle course. The course may be empty or use non-text content types.'}), 400

        return jsonify({
            'success': True,
            'slides': all_slides,
            'course_title': course_title,
            'slide_count': len(all_slides),
        })

    except Exception as e:
        print(f"[extract-moodle] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


@app.route('/extract-moodle-media', methods=['POST'])
def extract_moodle_media():
    """Find and transcribe media files from a Moodle course's SCORM packages."""
    try:
        import requests as http_requests
        import zipfile
        import io
        import uuid

        data = request.get_json()
        course_id = data.get('course_id')
        moodle_url = data.get('moodle_url')
        moodle_token = data.get('moodle_token')
        user_id = data.get('user_id')

        if not course_id or not moodle_url or not moodle_token or not user_id:
            return jsonify({'error': 'course_id, moodle_url, moodle_token, and user_id are required'}), 400

        base = moodle_url.rstrip('/')
        contents = _moodle_api_call(moodle_url, moodle_token, 'core_course_get_contents', {'courseid': course_id})

        MEDIA_EXTS = {'.mp3', '.wav', '.ogg', '.m4a', '.mp4', '.webm', '.mov', '.avi'}
        media_files = []

        for section in (contents if isinstance(contents, list) else []):
            for module in section.get('modules', []):
                if module.get('modname') == 'scorm':
                    for content_file in module.get('contents', []):
                        file_url = content_file.get('fileurl')
                        if file_url and file_url.endswith('.zip'):
                            try:
                                dl_url = f"{file_url}&token={moodle_token}" if '?' in file_url else f"{file_url}?token={moodle_token}"
                                resp = http_requests.get(dl_url, timeout=120)
                                if resp.status_code == 200:
                                    zf = zipfile.ZipFile(io.BytesIO(resp.content))
                                    for name in zf.namelist():
                                        ext = os.path.splitext(name)[1].lower()
                                        if ext in MEDIA_EXTS:
                                            media_files.append({
                                                'filename': name,
                                                'data': zf.read(name),
                                            })
                                    zf.close()
                            except Exception as e:
                                print(f"[extract-moodle-media] Warning: failed to process {file_url}: {e}")

        if not media_files:
            return jsonify({'success': True, 'job_id': None, 'media_files_found': 0})

        # Create async job and spawn transcription worker
        job_id = str(uuid.uuid4())
        supabase = get_supabase_client()
        supabase.table('generation_jobs').insert({
            'id': job_id,
            'user_id': user_id,
            'job_type': 'scorm_transcription',
            'status': 'processing',
        }).execute()

        import threading
        thread = threading.Thread(target=_transcription_job_worker, args=(job_id, media_files, course_id, user_id))
        thread.daemon = True
        thread.start()

        return jsonify({
            'success': True,
            'job_id': job_id,
            'media_files_found': len(media_files),
        })

    except Exception as e:
        print(f"[extract-moodle-media] Error: {e}")
        import traceback; traceback.print_exc()
        return jsonify({'error': str(e)}), 500


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port)
