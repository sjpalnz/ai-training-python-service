from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ── Built-in themes ────────────────────────────────────────────────────────────
# Each theme defines: slide_bg (None = white default), heading_color, text_color,
# accent_color (used for title-slide large text).  All values are RGBColor tuples.

THEMES = {
    'corporate': {
        'slide_bg':      None,                      # white
        'heading_color': RGBColor(30,  58,  95),    # deep navy
        'text_color':    RGBColor(51,  51,  51),    # dark charcoal
        'accent_color':  RGBColor(44,  90, 160),    # brand blue
        'title_size':    Pt(44),
        'heading_size':  Pt(32),
        'body_size':     Pt(18),
    },
    'dark': {
        'slide_bg':      RGBColor(18,  18,  32),    # near-black navy
        'heading_color': RGBColor(255, 255, 255),   # white
        'text_color':    RGBColor(200, 210, 230),   # light blue-grey
        'accent_color':  RGBColor(79, 172, 254),    # electric blue
        'title_size':    Pt(44),
        'heading_size':  Pt(32),
        'body_size':     Pt(18),
    },
    'slate': {
        'slide_bg':      RGBColor(245, 246, 248),   # very light grey
        'heading_color': RGBColor(45,  52,  54),    # charcoal
        'text_color':    RGBColor(45,  52,  54),    # charcoal
        'accent_color':  RGBColor(99, 159, 255),    # slate blue
        'title_size':    Pt(44),
        'heading_size':  Pt(32),
        'body_size':     Pt(18),
    },
    'clean': {
        'slide_bg':      RGBColor(255, 255, 255),   # pure white
        'heading_color': RGBColor(26,  26,  42),    # near-black
        'text_color':    RGBColor(60,  60,  80),    # dark purple-grey
        'accent_color':  RGBColor(48, 209,  88),    # green
        'title_size':    Pt(44),
        'heading_size':  Pt(32),
        'body_size':     Pt(18),
    },
    'vibrant': {
        'slide_bg':      RGBColor(255, 255, 255),   # white
        'heading_color': RGBColor(0,  119, 182),    # ocean blue
        'text_color':    RGBColor(51,  51,  51),    # dark grey
        'accent_color':  RGBColor(255, 107,  53),   # orange
        'title_size':    Pt(44),
        'heading_size':  Pt(32),
        'body_size':     Pt(18),
    },
}


# ── Helpers ────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, rgb_color):
    """Apply a solid background colour to a slide."""
    if rgb_color is None:
        return
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color


def clear_presentation_slides(prs):
    """
    Remove all existing slides from a loaded template presentation so we can
    add fresh content while keeping the slide master / theme intact.
    """
    xml_slides = prs.slides._sldIdLst
    slides = list(prs.slides)
    for slide in slides:
        slide_part = slide.part
        for rel in list(prs.part.rels.values()):
            if rel.target_part == slide_part:
                prs.part.drop_rel(rel.rId)
                break
    for element in list(xml_slides):
        xml_slides.remove(element)


def _set_font(run, size, bold=False, color=None):
    run.font.size = size
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_title_box(slide, text, left, top, width, height, size, bold, color, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = text
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            _set_font(run, size, bold=bold, color=color)
    return box


def _add_bullets(slide, bullets, left, top, width, height, size, color):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = size
        p.font.color.rgb = color


# ── Main generator ─────────────────────────────────────────────────────────────

def generate_powerpoint_file(course_data, output_path, theme_id='corporate', template_path=None):
    """
    Generate a PowerPoint from a course outline dict.

    Args:
        course_data:   Dict with 'title' and 'slides' array.
        output_path:   Destination .pptx path.
        theme_id:      One of 'corporate', 'dark', 'slate', 'clean', 'vibrant'.
        template_path: Optional path to a user-supplied .pptx template.  When
                       provided the template's slide master (background, fonts,
                       theme colours) is preserved; we skip setting a background
                       colour so the template's own design shows through.
    """
    theme = THEMES.get(theme_id, THEMES['corporate'])

    if template_path:
        prs = Presentation(template_path)
        clear_presentation_slides(prs)
        apply_bg = False   # let the template's master handle the background
    else:
        prs = Presentation()
        apply_bg = True

    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    heading_color = theme['heading_color']
    text_color    = theme['text_color']
    accent_color  = theme['accent_color']
    slide_bg      = theme['slide_bg'] if apply_bg else None

    title_size   = theme['title_size']
    heading_size = theme['heading_size']
    body_size    = theme['body_size']

    for slide_data in course_data.get('slides', []):
        slide_type = slide_data.get('type', 'content')
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
        set_slide_bg(slide, slide_bg)

        if slide_type == 'title':
            # Large centred title
            _add_title_box(
                slide,
                text   = slide_data.get('content', course_data.get('title', 'Training Course')),
                left   = Inches(0.5), top    = Inches(2.0),
                width  = Inches(9),   height = Inches(2),
                size   = title_size,  bold   = True,
                color  = heading_color, align = PP_ALIGN.CENTER,
            )
            # Sub-label (course title below main title if different)
            subtitle = course_data.get('title', '')
            if subtitle and subtitle != slide_data.get('content', ''):
                _add_title_box(
                    slide,
                    text   = subtitle,
                    left   = Inches(0.5), top    = Inches(4.2),
                    width  = Inches(9),   height = Inches(0.8),
                    size   = Pt(22),      bold   = False,
                    color  = accent_color, align = PP_ALIGN.CENTER,
                )

        elif slide_type == 'objectives':
            _add_title_box(
                slide,
                text   = slide_data.get('title', 'Learning Objectives'),
                left   = Inches(0.5), top    = Inches(0.5),
                width  = Inches(9),   height = Inches(0.9),
                size   = heading_size, bold  = True,
                color  = heading_color,
            )
            if slide_data.get('bullets'):
                _add_bullets(
                    slide, slide_data['bullets'],
                    left = Inches(1), top    = Inches(1.6),
                    width = Inches(8), height = Inches(5.5),
                    size = body_size, color  = text_color,
                )

        elif slide_type == 'content':
            _add_title_box(
                slide,
                text   = slide_data.get('title', 'Content'),
                left   = Inches(0.5), top    = Inches(0.5),
                width  = Inches(9),   height = Inches(0.9),
                size   = heading_size, bold  = True,
                color  = heading_color,
            )
            if slide_data.get('bullets'):
                _add_bullets(
                    slide, slide_data['bullets'],
                    left = Inches(1), top    = Inches(1.6),
                    width = Inches(8), height = Inches(5.5),
                    size = body_size, color  = text_color,
                )
            elif slide_data.get('content'):
                _add_title_box(
                    slide,
                    text   = slide_data['content'],
                    left   = Inches(0.5), top    = Inches(1.6),
                    width  = Inches(9),   height = Inches(5.5),
                    size   = body_size,   bold   = False,
                    color  = text_color,
                )

        elif slide_type == 'quiz':
            _add_title_box(
                slide,
                text   = slide_data.get('title', 'Knowledge Check'),
                left   = Inches(0.5), top    = Inches(0.4),
                width  = Inches(9),   height = Inches(0.8),
                size   = Pt(28),      bold   = True,
                color  = heading_color,
            )
            if slide_data.get('question'):
                _add_title_box(
                    slide,
                    text   = slide_data['question'],
                    left   = Inches(0.5), top    = Inches(1.3),
                    width  = Inches(9),   height = Inches(1.0),
                    size   = Pt(20),      bold   = True,
                    color  = text_color,
                )
            if slide_data.get('options'):
                _add_bullets(
                    slide, slide_data['options'],
                    left = Inches(1), top    = Inches(2.5),
                    width = Inches(8), height = Inches(4.5),
                    size = body_size, color  = text_color,
                )

        elif slide_type == 'summary':
            _add_title_box(
                slide,
                text   = slide_data.get('title', 'Key Takeaways'),
                left   = Inches(0.5), top    = Inches(0.5),
                width  = Inches(9),   height = Inches(0.9),
                size   = heading_size, bold  = True,
                color  = heading_color,
            )
            if slide_data.get('bullets'):
                _add_bullets(
                    slide, slide_data['bullets'],
                    left = Inches(1), top    = Inches(1.6),
                    width = Inches(8), height = Inches(5.5),
                    size = body_size, color  = text_color,
                )

        else:
            # Fallback for any unrecognised slide type
            _add_title_box(
                slide,
                text   = slide_data.get('title', ''),
                left   = Inches(0.5), top    = Inches(0.5),
                width  = Inches(9),   height = Inches(0.9),
                size   = heading_size, bold  = True,
                color  = heading_color,
            )
            if slide_data.get('content'):
                _add_title_box(
                    slide,
                    text   = slide_data['content'],
                    left   = Inches(0.5), top    = Inches(1.6),
                    width  = Inches(9),   height = Inches(5.5),
                    size   = body_size,   bold   = False,
                    color  = text_color,
                )

    prs.save(output_path)
    return output_path
